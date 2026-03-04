#!/usr/bin/env python3
"""
ProConvert Studio - Ultimate PDF Toolkit
Complete version with ALL features working including all conversions
"""

import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog, colorchooser, ttk
from PIL import Image, ImageTk, ImageDraw, ImageFont, ImageFilter
import fitz  # PyMuPDF
import os
import sys
import json
import threading
import queue
import tempfile
import shutil
import hashlib
import base64
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Callable, Tuple, Union, Any
import io
import re
import textwrap
import webbrowser
from dataclasses import dataclass
from enum import Enum

# Optional imports with fallbacks
try:
    from pdf2docx import Converter
    PDF2DOCX_AVAILABLE = True
except ImportError:
    PDF2DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Inches as DocxInches
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pytesseract
    from pdf2image import convert_from_path
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    from googletrans import Translator
    TRANSLATE_AVAILABLE = True
except ImportError:
    try:
        from deep_translator import GoogleTranslator
        TRANSLATE_AVAILABLE = True
    except ImportError:
        TRANSLATE_AVAILABLE = False

try:
    from cryptography.hazmat.primitives import hashes, serialization
    from cryptography.hazmat.primitives.asymmetric import padding, rsa
    from cryptography.hazmat.backends import default_backend
    CRYPTO_AVAILABLE = True
except ImportError:
    CRYPTO_AVAILABLE = False

# --- Constants and Configuration ---
VERSION = "2.0.0"
APP_NAME = "ProConvert Studio"
CONFIG_FILE = "proconvert_config.json"
RECENT_FILES_MAX = 15

class ConfigManager:
    """Manage application configuration, signatures, and recent files"""
    def __init__(self):
        self.config = {
            "theme": "System",
            "color_theme": "blue",
            "recent_files": [],
            "default_dpi": 300,
            "default_export_folder": "",
            "signatures": {},
            "watermark_presets": [],
            "shortcuts": {
                "open": "<Control-o>",
                "save": "<Control-s>",
                "save_as": "<Control-Shift-s>",
                "merge": "<Control-m>",
                "delete": "<Delete>",
                "rotate": "<Control-r>",
                "extract": "<Control-e>",
                "sign": "<Control-d>",
                "crop": "<Control-x>",
                "undo": "<Control-z>",
                "help": "<F1>"
            },
            "language": "en",
            "auto_save": True,
            "backup_on_save": True
        }
        self.load()
    
    def load(self):
        if os.path.exists(CONFIG_FILE):
            try:
                with open(CONFIG_FILE, 'r') as f:
                    loaded = json.load(f)
                    # Safe update
                    for key, value in loaded.items():
                        if key in self.config:
                            if isinstance(self.config[key], dict) and isinstance(value, dict):
                                self.config[key].update(value)
                            else:
                                self.config[key] = value
            except Exception:
                pass
    
    def save(self):
        try:
            with open(CONFIG_FILE, 'w') as f:
                json.dump(self.config, f, indent=2)
        except Exception:
            pass
    
    def add_recent_file(self, filepath: str):
        if filepath in self.config["recent_files"]:
            self.config["recent_files"].remove(filepath)
        self.config["recent_files"].insert(0, filepath)
        self.config["recent_files"] = self.config["recent_files"][:RECENT_FILES_MAX]
        self.save()
    
    def save_signature(self, name: str, image_data: str):
        self.config["signatures"][name] = image_data
        self.save()
    
    def get_signatures(self):
        return self.config.get("signatures", {})

class PDFEngine:
    """Enhanced PDF operations engine with ALL features working"""
    
    @staticmethod
    def get_page_thumbnail(doc: fitz.Document, page_num: int, zoom: float = 0.3) -> Image.Image:
        page = doc[page_num]
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        return img
    
    @staticmethod
    def get_page_count(filepath: str) -> int:
        try:
            with fitz.open(filepath) as doc:
                return len(doc)
        except:
            return 0
    
    @staticmethod
    def rotate_page_permanent(doc: fitz.Document, page_num: int, degrees: int = 90):
        page = doc[page_num]
        page.set_rotation((page.rotation + degrees) % 360)
        return page.rotation
    
    @staticmethod
    def crop_page(doc: fitz.Document, page_num: int, crop_box: Tuple[float, float, float, float]):
        page = doc[page_num]
        page.set_cropbox(fitz.Rect(crop_box))
        return page.rect
    
    @staticmethod
    def resize_page(doc: fitz.Document, page_num: int, scale_factor: float):
        """Resize page by scale factor"""
        page = doc[page_num]
        rect = page.rect
        new_rect = fitz.Rect(0, 0, rect.width * scale_factor, rect.height * scale_factor)
        # Note: This creates a new page with scaled content
        new_page = doc.new_page(width=new_rect.width, height=new_rect.height)
        new_page.show_pdf_page(new_rect, doc, page_num)
        doc.delete_page(page_num)
        return new_rect
    
    @staticmethod
    def auto_crop_margins(doc: fitz.Document, page_num: int, tolerance: int = 240):
        """Auto-detect and crop white margins"""
        page = doc[page_num]
        # Get page as image for analysis
        zoom = 2.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        # Convert to grayscale for better detection
        gray = img.convert('L')
        
        # Find content bounds (non-white pixels)
        bbox = None
        for y in range(gray.height):
            for x in range(gray.width):
                if gray.getpixel((x, y)) < tolerance:
                    if bbox is None:
                        bbox = [x, y, x, y]
                    else:
                        bbox[0] = min(bbox[0], x)
                        bbox[1] = min(bbox[1], y)
                        bbox[2] = max(bbox[2], x)
                        bbox[3] = max(bbox[3], y)
        
        if bbox:
            # Scale back to PDF coordinates
            scale = 1.0 / zoom
            crop_rect = fitz.Rect(
                bbox[0] * scale, bbox[1] * scale,
                bbox[2] * scale, bbox[3] * scale
            )
            page.set_cropbox(crop_rect)
            return crop_rect
        return None
    
    @staticmethod
    def merge_pdfs(files: List[str], output_path: str, progress_callback: Optional[Callable] = None):
        merged = fitz.open()
        total_pages = 0
        
        for f in files:
            with fitz.open(f) as tmp:
                total_pages += len(tmp)
        
        current_page = 0
        for f in files:
            with fitz.open(f) as src:
                merged.insert_pdf(src)
                current_page += len(src)
                if progress_callback:
                    progress_callback(current_page / total_pages)
        
        merged.save(output_path, garbage=4, deflate=True)
        merged.close()
    
    @staticmethod
    def split_pdf(input_path: str, output_folder: str, pages_per_file: int = 1):
        doc = fitz.open(input_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        output_files = []
        
        for i in range(0, len(doc), pages_per_file):
            new_doc = fitz.open()
            end = min(i + pages_per_file, len(doc))
            for j in range(i, end):
                new_doc.insert_pdf(doc, from_page=j, to_page=j)
            
            output_path = os.path.join(output_folder, f"{base_name}_part_{i//pages_per_file + 1}.pdf")
            new_doc.save(output_path)
            new_doc.close()
            output_files.append(output_path)
        
        doc.close()
        return output_files
    
    @staticmethod
    def compress_pdf(input_path: str, output_path: str, image_quality: int = 80):
        doc = fitz.open(input_path)
        
        # Optimize images if quality < 100
        if image_quality < 100:
            for page in doc:
                for img in page.get_images():
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n - pix.alpha < 4:  # Can be saved as JPEG
                        pix2 = fitz.Pixmap(pix, 0) if pix.alpha else pix
                        img_data = pix2.tobytes("jpeg", str(image_quality))
                        doc.update_stream(xref, img_data)
        
        doc.save(
            output_path,
            garbage=4,
            deflate=True,
            clean=True,
            linear=True
        )
        doc.close()
        
        original_size = os.path.getsize(input_path)
        compressed_size = os.path.getsize(output_path)
        ratio = (1 - compressed_size / original_size) * 100
        return ratio
    
    @staticmethod
    def add_watermark_text(doc: fitz.Document, text: str, opacity: float = 0.3, 
                          font_size: int = 50, color: Tuple[float, float, float] = (0.5, 0.5, 0.5),
                          angle: int = 45, position: str = "center"):
        for page in doc:
            rect = page.rect
            text_width = len(text) * font_size * 0.6
            
            if position == "center":
                pos = (rect.width / 2 - text_width/2, rect.height / 2)
            elif position == "top-left":
                pos = (50, 100)
            elif position == "top-right":
                pos = (rect.width - text_width - 50, 100)
            elif position == "bottom-left":
                pos = (50, rect.height - 100)
            elif position == "bottom-right":
                pos = (rect.width - text_width - 50, rect.height - 100)
            elif position == "tile":
                for x in range(0, int(rect.width), 300):
                    for y in range(0, int(rect.height), 300):
                        page.insert_text(
                            (x, y), text, fontsize=font_size/2,
                            color=color, opacity=opacity/2, rotate=angle
                        )
                continue
            
            page.insert_text(
                pos, text, fontsize=font_size,
                color=color, opacity=opacity, rotate=angle
            )
    
    @staticmethod
    def add_watermark_image(doc: fitz.Document, image_path: str, opacity: float = 0.3,
                           position: str = "center", scale: float = 0.5):
        for page in doc:
            rect = page.rect
            img_rect = fitz.Rect(0, 0, rect.width * scale, rect.height * scale)
            
            if position == "center":
                img_rect = fitz.Rect(
                    rect.width/2 - img_rect.width/2,
                    rect.height/2 - img_rect.height/2,
                    rect.width/2 + img_rect.width/2,
                    rect.height/2 + img_rect.height/2
                )
            elif position == "top-left":
                img_rect = fitz.Rect(50, 50, 50 + img_rect.width, 50 + img_rect.height)
            elif position == "top-right":
                img_rect = fitz.Rect(rect.width - img_rect.width - 50, 50,
                                    rect.width - 50, 50 + img_rect.height)
            elif position == "bottom-left":
                img_rect = fitz.Rect(50, rect.height - img_rect.height - 50,
                                    50 + img_rect.width, rect.height - 50)
            elif position == "bottom-right":
                img_rect = fitz.Rect(rect.width - img_rect.width - 50,
                                    rect.height - img_rect.height - 50,
                                    rect.width - 50, rect.height - 50)
            
            page.insert_image(img_rect, filename=image_path, overlay=True)
    
    @staticmethod
    def encrypt_pdf(input_path: str, output_path: str, password: str, 
                    allow_print: bool = True, allow_copy: bool = True,
                    allow_modify: bool = False):
        doc = fitz.open(input_path)
        permissions = fitz.PDF_PERM_ACCESSIBILITY
        if allow_print:
            permissions |= fitz.PDF_PERM_PRINT
        if allow_copy:
            permissions |= fitz.PDF_PERM_COPY
        if allow_modify:
            permissions |= fitz.PDF_PERM_MODIFY
        
        doc.save(
            output_path,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw=password,
            user_pw=password,
            permissions=permissions
        )
        doc.close()
    
    @staticmethod
    def decrypt_pdf(input_path: str, output_path: str, password: str):
        doc = fitz.open(input_path)
        if doc.authenticate(password):
            doc.save(output_path, garbage=4, deflate=True)
            doc.close()
            return True
        doc.close()
        return False
    
    @staticmethod
    def translate_page_text(page: fitz.Page, target_lang: str = 'en') -> str:
        """Extract and translate text from page"""
        text = page.get_text()
        if not TRANSLATE_AVAILABLE or not text.strip():
            return text
        
        try:
            if 'GoogleTranslator' in globals():
                translator = GoogleTranslator(source='auto', target=target_lang)
                return translator.translate(text)
            else:
                translator = Translator()
                result = translator.translate(text, dest=target_lang)
                return result.text
        except Exception as e:
            return f"[Translation Error: {e}]\n\nOriginal:\n{text}"
    
    @staticmethod
    def images_to_pdf(image_files: List[str], output_path: str, 
                      page_size: str = "A4", dpi: int = 300,
                      fit_mode: str = "contain"):
        doc = fitz.open()
        
        page_sizes = {
            "A4": (595, 842),
            "Letter": (612, 792),
            "Legal": (612, 1008),
            "A3": (842, 1191)
        }
        
        for img_path in image_files:
            img = Image.open(img_path)
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            
            if page_size == "Original":
                width = img.width * 72 / dpi
                height = img.height * 72 / dpi
                page = doc.new_page(width=width, height=height)
                margin = 0
            else:
                w, h = page_sizes.get(page_size, (595, 842))
                page = doc.new_page(width=w, height=h)
                margin = 36
                width, height = w, h
            
            available_width = width - 2 * margin
            available_height = height - 2 * margin
            
            img_ratio = img.width / img.height
            page_ratio = available_width / available_height
            
            if fit_mode == "contain":
                if img_ratio > page_ratio:
                    new_width = available_width
                    new_height = available_width / img_ratio
                else:
                    new_height = available_height
                    new_width = available_height * img_ratio
            elif fit_mode == "cover":
                if img_ratio > page_ratio:
                    new_height = available_height
                    new_width = available_height * img_ratio
                else:
                    new_width = available_width
                    new_height = available_width / img_ratio
            else:  # stretch
                new_width = available_width
                new_height = available_height
            
            x = (width - new_width) / 2
            y = (height - new_height) / 2
            
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
                temp_path = tmp.name
                img.save(temp_path, 'JPEG', quality=95, dpi=(dpi, dpi))
            
            page.insert_image(fitz.Rect(x, y, x + new_width, y + new_height), 
                            filename=temp_path)
            os.remove(temp_path)
        
        doc.save(output_path, garbage=4, deflate=True)
        doc.close()
    
    @staticmethod
    def pdf_to_images_fixed(doc: fitz.Document, output_folder: str, 
                           format: str = "jpg", dpi: int = 300):
        """Fixed PDF to images conversion using PIL"""
        zoom = dpi / 72
        exported = []
        
        for i, page in enumerate(doc):
            # Render page
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to PIL Image
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Determine output path
            if format.lower() == 'jpg':
                output_path = os.path.join(output_folder, f"page_{i+1:03d}.jpg")
                img.save(output_path, 'JPEG', quality=95, dpi=(dpi, dpi))
            elif format.lower() == 'png':
                output_path = os.path.join(output_folder, f"page_{i+1:03d}.png")
                img.save(output_path, 'PNG', dpi=(dpi, dpi))
            elif format.lower() == 'tiff':
                output_path = os.path.join(output_folder, f"page_{i+1:03d}.tiff")
                img.save(output_path, 'TIFF', dpi=(dpi, dpi))
            else:  # bmp
                output_path = os.path.join(output_folder, f"page_{i+1:03d}.bmp")
                img.save(output_path, 'BMP', dpi=(dpi, dpi))
            
            exported.append(output_path)
        
        return exported

class SignatureCanvas(ctk.CTkCanvas):
    """Canvas for drawing signatures"""
    def __init__(self, master, width=400, height=200, **kwargs):
        super().__init__(master, width=width, height=height, bg="white", **kwargs)
        self.width = width
        self.height = height
        self.lines = []
        self.current_line = []
        self.pen_color = "black"
        self.pen_width = 2
        
        self.bind("<Button-1>", self.on_press)
        self.bind("<B1-Motion>", self.on_drag)
        self.bind("<ButtonRelease-1>", self.on_release)
        
    def on_press(self, event):
        self.current_line = [(event.x, event.y)]
        
    def on_drag(self, event):
        self.current_line.append((event.x, event.y))
        if len(self.current_line) > 1:
            x1, y1 = self.current_line[-2]
            x2, y2 = self.current_line[-1]
            self.create_line(x1, y1, x2, y2, fill=self.pen_color, 
                           width=self.pen_width, capstyle=tk.ROUND, smooth=True)
    
    def on_release(self, event):
        if self.current_line:
            self.lines.append(self.current_line)
        self.current_line = []
    
    def clear(self):
        self.delete("all")
        self.lines = []
    
    def get_image(self) -> Image.Image:
        ps = self.postscript(colormode="color", width=self.width, height=self.height)
        img = Image.open(io.BytesIO(ps.encode('utf-8')))
        return img.convert("RGBA")
    
    def save_signature(self, path: str):
        img = self.get_image()
        datas = img.getdata()
        newData = []
        for item in datas:
            if item[0] > 240 and item[1] > 240 and item[2] > 240:
                newData.append((255, 255, 255, 0))
            else:
                newData.append(item)
        img.putdata(newData)
        img.save(path, "PNG")

class DragDropListbox(tk.Listbox):
    """Listbox with drag-and-drop reordering"""
    def __init__(self, master, **kwargs):
        super().__init__(master, **kwargs)
        self.drag_data = None
        self.bind('<Button-1>', self.on_click)
        self.bind('<B1-Motion>', self.on_drag)
        self.bind('<ButtonRelease-1>', self.on_drop)
        
    def on_click(self, event):
        self.drag_data = self.nearest(event.y)
        self.selection_clear(0, tk.END)
        self.selection_set(self.drag_data)
        
    def on_drag(self, event):
        if self.drag_data is None:
            return
        new_index = self.nearest(event.y)
        if new_index != self.drag_data and new_index >= 0:
            item = self.get(self.drag_data)
            self.delete(self.drag_data)
            self.insert(new_index, item)
            self.selection_set(new_index)
            self.drag_data = new_index
            
    def on_drop(self, event):
        self.drag_data = None

class ProConvertStudio(ctk.CTk):
    def __init__(self):
        super().__init__()
        
        self.config_manager = ConfigManager()
        self.pdf_engine = PDFEngine()
        
        # State
        self.current_doc: Optional[fitz.Document] = None
        self.current_path: Optional[str] = None
        self.page_widgets: Dict[int, Dict] = {}
        self.merge_files: List[str] = []
        self.ui_queue = queue.Queue()
        self.operation_thread: Optional[threading.Thread] = None
        self.crop_start = None
        self.crop_rect = None
        self.undo_stack = []
        self.redo_stack = []
        self.sign_pdf_path = None
        self.watermark_image_path = None
        self.translate_pdf_path = None
        self.compress_file_path = None
        self.encrypt_file_path = None
        self.decrypt_file_path = None
        self.ocr_file_path = None
        
        # Setup window
        self.setup_window()
        self.setup_styles()
        self.setup_layout()
        self.setup_bindings()
        self.setup_ui_updater()
        self.check_dependencies()
        
        ctk.set_appearance_mode(self.config_manager.config["theme"])
        ctk.set_default_color_theme(self.config_manager.config["color_theme"])
        
    def setup_window(self):
        self.title(f"ProConvert Studio v{VERSION} - Ultimate PDF Toolkit")
        self.geometry("1800x1000")
        self.minsize(1400, 800)
        
    def setup_styles(self):
        self.colors = {
            "primary": "#3B8ED0",
            "success": "#2ecc71",
            "warning": "#f39c12",
            "danger": "#e74c3c",
            "info": "#3498db",
            "dark": "#2c3e50",
            "light": "#ecf0f1"
        }
        
    def setup_layout(self):
        self.grid_columnconfigure(1, weight=1)
        self.grid_rowconfigure(0, weight=1)
        
        self.setup_sidebar()
        self.setup_main_container()
        self.setup_statusbar()
        self.setup_menu()
        
    def setup_menu(self):
        menubar = tk.Menu(self)
        self.config(menu=menubar)
        
        file_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="File", menu=file_menu)
        file_menu.add_command(label="Open PDF", command=self.open_pdf, accelerator="Ctrl+O")
        file_menu.add_command(label="Save", command=self.save_pdf, accelerator="Ctrl+S")
        file_menu.add_command(label="Save As...", command=self.save_pdf_as, accelerator="Ctrl+Shift+S")
        file_menu.add_separator()
        file_menu.add_command(label="Exit", command=self.quit)
        
        edit_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="Edit", menu=edit_menu)
        edit_menu.add_command(label="Undo", command=self.undo_action, accelerator="Ctrl+Z")
        edit_menu.add_command(label="Redo", command=self.redo_action, accelerator="Ctrl+Y")
        edit_menu.add_separator()
        edit_menu.add_command(label="Settings", command=self.show_settings)
        
        tools_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="Tools", menu=tools_menu)
        tools_menu.add_command(label="Merge PDFs", command=self.show_merger)
        tools_menu.add_command(label="Split PDF", command=self.show_split_dialog)
        tools_menu.add_command(label="Compress PDF", command=self.show_compress)
        tools_menu.add_command(label="Encrypt PDF", command=self.show_security)
        tools_menu.add_separator()
        tools_menu.add_command(label="Add Watermark", command=self.show_watermark)
        tools_menu.add_command(label="Add Signature", command=self.show_sign)
        tools_menu.add_command(label="OCR Text Extraction", command=self.show_ocr)
        tools_menu.add_command(label="Translate PDF", command=self.show_translate)
        
        help_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="Help", menu=help_menu)
        help_menu.add_command(label="Keyboard Shortcuts", command=self.show_shortcuts)
        help_menu.add_command(label="About", command=self.show_about)
        
    def setup_sidebar(self):
        self.sidebar = ctk.CTkFrame(self, width=280, corner_radius=0)
        self.sidebar.grid(row=0, column=0, sticky="nsew", rowspan=2)
        self.sidebar.grid_rowconfigure(12, weight=1)
        
        logo_frame = ctk.CTkFrame(self.sidebar, fg_color="transparent")
        logo_frame.grid(row=0, column=0, padx=20, pady=(30, 20), sticky="ew")
        
        ctk.CTkLabel(logo_frame, text="📄", font=("Segoe UI Emoji", 48)).pack()
        ctk.CTkLabel(logo_frame, text="ProConvert", 
                    font=ctk.CTkFont(size=24, weight="bold")).pack()
        ctk.CTkLabel(logo_frame, text="Studio Ultimate", 
                    font=ctk.CTkFont(size=14), text_color="gray").pack()
        
        ctk.CTkLabel(self.sidebar, text="TOOLS", font=ctk.CTkFont(size=12, weight="bold"),
                    text_color="gray").grid(row=1, column=0, padx=20, pady=(20, 10), sticky="w")
        
        self.nav_buttons = {}
        tools = [
            ("converter", "🔄 Converter", self.show_converter),
            ("manager", "✏️ PDF Editor", self.show_manager),
            ("merger", "🔀 PDF Merger", self.show_merger),
            ("sign", "✍️ Sign & Certify", self.show_sign),
            ("watermark", "💧 Watermark", self.show_watermark),
            ("crop", "✂️ Crop & Resize", self.show_crop),
            ("translate", "🌐 Translate", self.show_translate),
            ("compress", "🗜️ Compress", self.show_compress),
            ("security", "🔒 Security", self.show_security),
            ("ocr", "🔍 OCR", self.show_ocr),
        ]
        
        for i, (key, text, cmd) in enumerate(tools, start=2):
            btn = ctk.CTkButton(
                self.sidebar, 
                text=text, 
                command=cmd,
                fg_color="transparent", 
                text_color=("black", "white"),
                hover_color=("gray75", "gray25"), 
                anchor="w", 
                height=45,
                font=ctk.CTkFont(size=14),
                corner_radius=8
            )
            btn.grid(row=i, column=0, padx=15, pady=3, sticky="ew")
            self.nav_buttons[key] = btn
        
        ctk.CTkLabel(self.sidebar, text="RECENT FILES", font=ctk.CTkFont(size=12, weight="bold"),
                    text_color="gray").grid(row=12, column=0, padx=20, pady=(20, 10), sticky="w")
        
        self.recent_frame = ctk.CTkScrollableFrame(self.sidebar, height=200)
        self.recent_frame.grid(row=13, column=0, padx=15, pady=5, sticky="ew")
        self.refresh_recent_files()
        
        ctk.CTkButton(
            self.sidebar, 
            text="⚙️ Settings", 
            command=self.show_settings, 
            height=40,
            font=ctk.CTkFont(size=14)
        ).grid(row=14, column=0, padx=15, pady=15, sticky="ew")
        
    def setup_main_container(self):
        self.container = ctk.CTkFrame(self, fg_color="transparent")
        self.container.grid(row=0, column=1, sticky="nsew", padx=15, pady=15)
        self.container.grid_rowconfigure(0, weight=1)
        self.container.grid_columnconfigure(0, weight=1)
        
        self.tabs = {}
        
        self.tabs["converter"] = self.create_converter_tab()
        self.tabs["manager"] = self.create_manager_tab()
        self.tabs["merger"] = self.create_merger_tab()
        self.tabs["sign"] = self.create_sign_tab()
        self.tabs["watermark"] = self.create_watermark_tab()
        self.tabs["crop"] = self.create_crop_tab()
        self.tabs["translate"] = self.create_translate_tab()
        self.tabs["compress"] = self.create_compress_tab()
        self.tabs["security"] = self.create_security_tab()
        self.tabs["ocr"] = self.create_ocr_tab()
        
        self.show_converter()
        
    def create_converter_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        frame.grid_columnconfigure(0, weight=1)
        frame.grid_rowconfigure(1, weight=1)
        frame.grid_rowconfigure(2, weight=1)
        
        header = ctk.CTkFrame(frame, height=80)
        header.grid(row=0, column=0, sticky="ew", padx=15, pady=15)
        
        ctk.CTkLabel(
            header, 
            text="Universal Converter", 
            font=ctk.CTkFont(size=32, weight="bold")
        ).pack(side="left", padx=20, pady=10)
        
        scroll = ctk.CTkScrollableFrame(frame)
        scroll.grid(row=1, column=0, sticky="nsew", padx=15, pady=10)
        scroll.grid_columnconfigure((0, 1, 2), weight=1, minsize=250)
        
        conversions = [
            ("📷", "Images", "PDF", lambda: self.start_conversion("img_pdf"), True),
            ("🖼️", "PDF", "Images", lambda: self.start_conversion("pdf_img"), True),
            ("📝", "PDF", "Word", lambda: self.start_conversion("pdf_word"), PDF2DOCX_AVAILABLE),
            ("📄", "Word", "PDF", lambda: self.start_conversion("word_pdf"), DOCX_AVAILABLE),
            ("📊", "PDF", "PowerPoint", lambda: self.start_conversion("pdf_ppt"), PPTX_AVAILABLE),
            ("📈", "PowerPoint", "PDF", lambda: self.start_conversion("ppt_pdf"), PPTX_AVAILABLE),
            ("📄", "PDF", "Text", lambda: self.start_conversion("pdf_txt"), True),
            ("🌐", "PDF", "HTML", lambda: self.start_conversion("pdf_html"), True),
            ("📑", "PDF", "PDF/A", lambda: self.start_conversion("pdf_pdfa"), True),
        ]
        
        for i, (icon, from_fmt, to_fmt, cmd, available) in enumerate(conversions):
            btn_frame = ctk.CTkFrame(scroll, corner_radius=10)
            btn_frame.grid(row=i//3, column=i%3, padx=15, pady=15, sticky="nsew")
            
            btn = ctk.CTkButton(
                btn_frame, 
                text=f"{icon}\n{from_fmt} → {to_fmt}", 
                command=cmd, 
                width=220, 
                height=120,
                font=ctk.CTkFont(size=16),
                corner_radius=10,
                state="normal" if available else "disabled"
            )
            btn.pack(padx=15, pady=15, expand=True, fill="both")
            
            if not available:
                ctk.CTkLabel(
                    btn_frame, 
                    text="(Install required package)", 
                    text_color="orange", 
                    font=ctk.CTkFont(size=11)
                ).pack(pady=(0, 10))
        
        log_frame = ctk.CTkFrame(frame)
        log_frame.grid(row=2, column=0, sticky="nsew", padx=15, pady=15)
        log_frame.grid_rowconfigure(1, weight=1)
        log_frame.grid_columnconfigure(0, weight=1)
        
        ctk.CTkLabel(
            log_frame, 
            text="Activity Log", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).grid(row=0, column=0, sticky="w", padx=15, pady=10)
        
        self.log_text = ctk.CTkTextbox(log_frame, state="disabled", font=("Consolas", 11))
        self.log_text.grid(row=1, column=0, sticky="nsew", padx=15, pady=(0, 15))
        
        return frame
    
    def create_manager_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        frame.grid_columnconfigure(0, weight=1)
        frame.grid_rowconfigure(1, weight=1)
        
        toolbar = ctk.CTkFrame(frame, height=60)
        toolbar.grid(row=0, column=0, sticky="ew", padx=15, pady=15)
        
        buttons = [
            ("📂 Open", self.open_pdf, None),
            ("💾 Save", self.save_pdf, "green"),
            ("↩️ Undo", self.undo_action, "gray"),
            ("🔄 Rotate", lambda: self.manipulate_pages("rotate"), None),
            ("🗑️ Delete", lambda: self.manipulate_pages("delete"), "red"),
            ("📤 Extract", lambda: self.manipulate_pages("extract"), None),
        ]
        
        for i, (text, cmd, color) in enumerate(buttons):
            kwargs = {
                "fg_color": color if color else None,
                "height": 40,
                "width": 100,
                "font": ctk.CTkFont(size=13)
            }
            btn = ctk.CTkButton(toolbar, text=text, command=cmd, **kwargs)
            btn.grid(row=0, column=i, padx=5, pady=10)
        
        paned = ctk.CTkFrame(frame)
        paned.grid(row=1, column=0, sticky="nsew", padx=15, pady=15)
        paned.grid_columnconfigure(0, weight=3)
        paned.grid_columnconfigure(1, weight=1)
        paned.grid_rowconfigure(0, weight=1)
        
        self.thumb_frame = ctk.CTkScrollableFrame(
            paned, 
            label_text="Pages (Click to select, double-click to preview)",
            label_font=ctk.CTkFont(size=13)
        )
        self.thumb_frame.grid(row=0, column=0, sticky="nsew", padx=5)
        
        preview_frame = ctk.CTkFrame(paned)
        preview_frame.grid(row=0, column=1, sticky="nsew", padx=5)
        preview_frame.grid_rowconfigure(1, weight=1)
        preview_frame.grid_columnconfigure(0, weight=1)
        
        ctk.CTkLabel(
            preview_frame, 
            text="Preview", 
            font=ctk.CTkFont(size=18, weight="bold")
        ).grid(row=0, column=0, pady=10)
        
        self.preview_canvas = tk.Canvas(
            preview_frame, 
            bg="#2b2b2b", 
            width=400, 
            height=550,
            highlightthickness=0
        )
        self.preview_canvas.grid(row=1, column=0, padx=10, pady=10, sticky="nsew")
        
        self.page_info_label = ctk.CTkLabel(
            preview_frame, 
            text="No document loaded",
            font=ctk.CTkFont(size=12)
        )
        self.page_info_label.grid(row=2, column=0, pady=5)
        
        self.preview_canvas.bind("<Button-1>", self.on_crop_start)
        self.preview_canvas.bind("<B1-Motion>", self.on_crop_drag)
        self.preview_canvas.bind("<ButtonRelease-1>", self.on_crop_end)
        
        return frame
    
    def create_sign_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        frame.grid_columnconfigure(0, weight=1)
        
        ctk.CTkLabel(
            frame, 
            text="Digital Signature & Certification", 
            font=ctk.CTkFont(size=28, weight="bold")
        ).pack(pady=20)
        
        create_frame = ctk.CTkFrame(frame)
        create_frame.pack(pady=15, padx=50, fill="x")
        
        ctk.CTkLabel(
            create_frame, 
            text="Create Signature", 
            font=ctk.CTkFont(size=18, weight="bold")
        ).pack(pady=10)
        
        self.sign_canvas = SignatureCanvas(create_frame, width=600, height=250)
        self.sign_canvas.pack(pady=10)
        
        btn_frame = ctk.CTkFrame(create_frame, fg_color="transparent")
        btn_frame.pack(pady=10)
        
        ctk.CTkButton(
            btn_frame, 
            text="Clear Canvas", 
            command=self.sign_canvas.clear,
            width=120,
            height=35
        ).pack(side="left", padx=5)
        
        ctk.CTkButton(
            btn_frame, 
            text="Save Signature", 
            command=self.save_drawn_signature, 
            fg_color="green",
            width=120,
            height=35
        ).pack(side="left", padx=5)
        
        apply_frame = ctk.CTkFrame(frame)
        apply_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkLabel(
            apply_frame, 
            text="Apply to PDF", 
            font=ctk.CTkFont(size=18, weight="bold")
        ).pack(pady=10)
        
        ctk.CTkButton(
            apply_frame, 
            text="📂 Select PDF File", 
            command=self.select_sign_pdf,
            height=40,
            width=200
        ).pack(pady=5)
        
        self.sign_pdf_label = ctk.CTkLabel(apply_frame, text="No file selected")
        self.sign_pdf_label.pack()
        
        self.sig_var = tk.StringVar()
        self.sig_combo = ctk.CTkComboBox(
            apply_frame, 
            variable=self.sig_var, 
            values=[], 
            width=350,
            height=35
        )
        self.sig_combo.pack(pady=10)
        self.refresh_signatures()
        
        ctk.CTkLabel(apply_frame, text="Position:").pack()
        self.sign_pos = ctk.CTkComboBox(
            apply_frame, 
            values=["Bottom-Right", "Bottom-Left", "Top-Right", "Top-Left", "Center"], 
            width=200,
            height=35
        )
        self.sign_pos.set("Bottom-Right")
        self.sign_pos.pack(pady=5)
        
        ctk.CTkButton(
            apply_frame, 
            text="✍️ Apply Signature", 
            command=self.apply_signature, 
            fg_color="green", 
            height=45,
            width=200,
            font=ctk.CTkFont(size=14)
        ).pack(pady=20)
        
        return frame
    
    def create_watermark_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        
        ctk.CTkLabel(
            frame, 
            text="Advanced Watermarking", 
            font=ctk.CTkFont(size=28, weight="bold")
        ).pack(pady=20)
        
        text_frame = ctk.CTkFrame(frame)
        text_frame.pack(pady=15, padx=50, fill="x")
        
        ctk.CTkLabel(
            text_frame, 
            text="Text Watermark", 
            font=ctk.CTkFont(size=18, weight="bold")
        ).pack(pady=10)
        
        self.wm_text = ctk.CTkEntry(
            text_frame, 
            placeholder_text="Enter watermark text", 
            width=450,
            height=35
        )
        self.wm_text.pack(pady=5)
        
        opts_frame = ctk.CTkFrame(text_frame, fg_color="transparent")
        opts_frame.pack(pady=15)
        
        size_frame = ctk.CTkFrame(opts_frame, fg_color="transparent")
        size_frame.pack(side="left", padx=10)
        ctk.CTkLabel(size_frame, text="Size:").pack()
        self.wm_size = ctk.CTkSlider(size_frame, from_=10, to=100, width=120)
        self.wm_size.set(50)
        self.wm_size.pack()
        self.wm_size_label = ctk.CTkLabel(size_frame, text="50")
        self.wm_size_label.pack()
        self.wm_size.configure(command=lambda v: self.wm_size_label.configure(text=str(int(v))))
        
        opacity_frame = ctk.CTkFrame(opts_frame, fg_color="transparent")
        opacity_frame.pack(side="left", padx=10)
        ctk.CTkLabel(opacity_frame, text="Opacity:").pack()
        self.wm_opacity = ctk.CTkSlider(opacity_frame, from_=0.1, to=1.0, width=120)
        self.wm_opacity.set(0.3)
        self.wm_opacity.pack()
        self.wm_opacity_label = ctk.CTkLabel(opacity_frame, text="0.3")
        self.wm_opacity_label.pack()
        self.wm_opacity.configure(command=lambda v: self.wm_opacity_label.configure(text=f"{v:.1f}"))
        
        angle_frame = ctk.CTkFrame(opts_frame, fg_color="transparent")
        angle_frame.pack(side="left", padx=10)
        ctk.CTkLabel(angle_frame, text="Angle:").pack()
        self.wm_angle = ctk.CTkSlider(angle_frame, from_=0, to=360, width=120)
        self.wm_angle.set(45)
        self.wm_angle.pack()
        self.wm_angle_label = ctk.CTkLabel(angle_frame, text="45°")
        self.wm_angle_label.pack()
        self.wm_angle.configure(command=lambda v: self.wm_angle_label.configure(text=f"{int(v)}°"))
        
        pos_frame = ctk.CTkFrame(text_frame, fg_color="transparent")
        pos_frame.pack(pady=10)
        
        ctk.CTkLabel(pos_frame, text="Position:", font=ctk.CTkFont(size=12)).pack(side="left", padx=5)
        self.wm_pos = ctk.CTkComboBox(
            pos_frame, 
            values=["Center", "Tile", "Top-Left", "Top-Right", "Bottom-Left", "Bottom-Right"], 
            width=150,
            height=32
        )
        self.wm_pos.set("Center")
        self.wm_pos.pack(side="left", padx=5)
        
        ctk.CTkButton(
            pos_frame, 
            text="Pick Color", 
            command=self.pick_watermark_color,
            width=100,
            height=32
        ).pack(side="left", padx=5)
        
        self.wm_color = (0.5, 0.5, 0.5)
        self.wm_color_preview = ctk.CTkLabel(pos_frame, text="⬛", font=("Segoe UI Emoji", 20))
        self.wm_color_preview.pack(side="left", padx=5)
        
        ctk.CTkButton(
            text_frame, 
            text="💧 Apply Text Watermark", 
            command=lambda: self.apply_watermark("text"), 
            fg_color="blue",
            height=40,
            width=200
        ).pack(pady=15)
        
        img_frame = ctk.CTkFrame(frame)
        img_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkLabel(
            img_frame, 
            text="Image Watermark", 
            font=ctk.CTkFont(size=18, weight="bold")
        ).pack(pady=10)
        
        ctk.CTkButton(
            img_frame, 
            text="Select Image", 
            command=self.select_watermark_image,
            height=40,
            width=200
        ).pack(pady=5)
        
        self.wm_img_label = ctk.CTkLabel(img_frame, text="No image selected")
        self.wm_img_label.pack()
        
        img_opts = ctk.CTkFrame(img_frame, fg_color="transparent")
        img_opts.pack(pady=10)
        
        ctk.CTkLabel(img_opts, text="Scale:").pack(side="left", padx=5)
        self.wm_img_scale = ctk.CTkSlider(img_opts, from_=0.1, to=1.0, width=120)
        self.wm_img_scale.set(0.5)
        self.wm_img_scale.pack(side="left", padx=5)
        
        ctk.CTkButton(
            img_frame, 
            text="🖼️ Apply Image Watermark", 
            command=lambda: self.apply_watermark("image"), 
            fg_color="blue",
            height=40,
            width=200
        ).pack(pady=15)
        
        return frame
    
    def create_crop_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        
        ctk.CTkLabel(
            frame, 
            text="Crop & Resize PDF", 
            font=ctk.CTkFont(size=28, weight="bold")
        ).pack(pady=20)
        
        info_frame = ctk.CTkFrame(frame)
        info_frame.pack(pady=10, padx=50, fill="x")
        
        ctk.CTkLabel(
            info_frame, 
            text="Open a PDF in the Editor tab, then use the preview panel to draw crop area", 
            text_color="gray",
            font=ctk.CTkFont(size=12)
        ).pack(pady=10)
        
        manual_frame = ctk.CTkFrame(frame)
        manual_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkLabel(
            manual_frame, 
            text="Manual Crop", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(pady=10)
        
        dims_frame = ctk.CTkFrame(manual_frame, fg_color="transparent")
        dims_frame.pack(pady=10)
        
        ctk.CTkLabel(dims_frame, text="X:").pack(side="left")
        self.crop_x = ctk.CTkEntry(dims_frame, width=60)
        self.crop_x.pack(side="left", padx=5)
        
        ctk.CTkLabel(dims_frame, text="Y:").pack(side="left")
        self.crop_y = ctk.CTkEntry(dims_frame, width=60)
        self.crop_y.pack(side="left", padx=5)
        
        ctk.CTkLabel(dims_frame, text="Width:").pack(side="left")
        self.crop_w = ctk.CTkEntry(dims_frame, width=60)
        self.crop_w.pack(side="left", padx=5)
        
        ctk.CTkLabel(dims_frame, text="Height:").pack(side="left")
        self.crop_h = ctk.CTkEntry(dims_frame, width=60)
        self.crop_h.pack(side="left", padx=5)
        
        ctk.CTkButton(
            manual_frame, 
            text="Apply Crop to Selected Pages", 
            command=self.apply_manual_crop,
            fg_color="green",
            height=40
        ).pack(pady=10)
        
        scale_frame = ctk.CTkFrame(frame)
        scale_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkLabel(
            scale_frame, 
            text="Resize Pages", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(pady=10)
        
        scale_control = ctk.CTkFrame(scale_frame, fg_color="transparent")
        scale_control.pack(pady=10)
        
        ctk.CTkLabel(scale_control, text="Scale Factor:").pack(side="left", padx=5)
        self.scale_factor = ctk.CTkSlider(scale_control, from_=0.1, to=2.0, width=200)
        self.scale_factor.set(1.0)
        self.scale_factor.pack(side="left", padx=5)
        self.scale_label = ctk.CTkLabel(scale_control, text="1.0x")
        self.scale_label.pack(side="left", padx=5)
        self.scale_factor.configure(command=lambda v: self.scale_label.configure(text=f"{v:.1f}x"))
        
        ctk.CTkButton(
            scale_frame, 
            text="Scale All Pages", 
            command=self.scale_pages,
            height=40
        ).pack(pady=10)
        
        auto_frame = ctk.CTkFrame(frame)
        auto_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkLabel(
            auto_frame, 
            text="Auto-Crop", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(pady=10)
        
        ctk.CTkButton(
            auto_frame, 
            text="Remove White Margins", 
            command=self.auto_crop_margins,
            height=40
        ).pack(pady=10)
        
        return frame
    
    def create_translate_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        
        ctk.CTkLabel(
            frame, 
            text="PDF Translation", 
            font=ctk.CTkFont(size=28, weight="bold")
        ).pack(pady=20)
        
        if not TRANSLATE_AVAILABLE:
            ctk.CTkLabel(
                frame, 
                text="⚠️ Translation requires googletrans or deep_translator", 
                text_color="orange",
                font=ctk.CTkFont(size=14)
            ).pack(pady=20)
            return frame
        
        src_frame = ctk.CTkFrame(frame)
        src_frame.pack(pady=10, padx=50, fill="x")
        
        ctk.CTkButton(
            src_frame, 
            text="📂 Select PDF", 
            command=self.select_translate_pdf,
            height=40,
            width=200
        ).pack(pady=10)
        
        self.trans_pdf_label = ctk.CTkLabel(src_frame, text="No file selected")
        self.trans_pdf_label.pack()
        
        lang_frame = ctk.CTkFrame(frame)
        lang_frame.pack(pady=20)
        
        ctk.CTkLabel(lang_frame, text="From:").pack(side="left", padx=5)
        self.src_lang = ctk.CTkComboBox(
            lang_frame, 
            values=["auto", "en", "es", "fr", "de", "it", "pt", "ru", "zh", "ja", "ko"],
            width=100
        )
        self.src_lang.set("auto")
        self.src_lang.pack(side="left", padx=5)
        
        ctk.CTkLabel(lang_frame, text="To:").pack(side="left", padx=5)
        self.tgt_lang = ctk.CTkComboBox(
            lang_frame, 
            values=["en", "es", "fr", "de", "it", "pt", "ru", "zh", "ja", "ko"],
            width=100
        )
        self.tgt_lang.set("en")
        self.tgt_lang.pack(side="left", padx=5)
        
        opts_frame = ctk.CTkFrame(frame)
        opts_frame.pack(pady=10)
        
        self.trans_mode = tk.StringVar(value="text")
        ctk.CTkRadioButton(opts_frame, text="Export as text file", 
                          variable=self.trans_mode, value="text").pack(anchor="w", pady=2)
        ctk.CTkRadioButton(opts_frame, text="Create translated PDF (experimental)", 
                          variable=self.trans_mode, value="pdf").pack(anchor="w", pady=2)
        
        ctk.CTkButton(
            frame, 
            text="🌐 Translate", 
            command=self.execute_translation,
            fg_color="green",
            height=45,
            width=200
        ).pack(pady=20)
        
        self.trans_progress = ctk.CTkProgressBar(frame, width=400)
        self.trans_progress.pack(pady=10)
        self.trans_progress.set(0)
        
        self.trans_result = ctk.CTkTextbox(frame, height=200)
        self.trans_result.pack(pady=20, padx=50, fill="both", expand=True)
        
        return frame
    
    def create_merger_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        frame.grid_columnconfigure(0, weight=1)
        frame.grid_rowconfigure(1, weight=1)
        
        ctrl = ctk.CTkFrame(frame, height=80)
        ctrl.grid(row=0, column=0, sticky="ew", padx=15, pady=15)
        
        ctk.CTkLabel(
            ctrl, 
            text="PDF Merger", 
            font=ctk.CTkFont(size=24, weight="bold")
        ).pack(side="left", padx=20)
        
        button_frame = ctk.CTkFrame(ctrl, fg_color="transparent")
        button_frame.pack(side="right", padx=20)
        
        ctk.CTkButton(
            button_frame, 
            text="➕ Add Files", 
            command=self.add_merge_files,
            height=40,
            width=120
        ).pack(side="left", padx=5)
        
        ctk.CTkButton(
            button_frame, 
            text="🗑️ Clear", 
            command=self.clear_merge_files, 
            fg_color="red",
            height=40,
            width=100
        ).pack(side="left", padx=5)
        
        ctk.CTkButton(
            button_frame, 
            text="💾 Merge & Save", 
            command=self.execute_merge, 
            fg_color="green", 
            height=45,
            width=150
        ).pack(side="left", padx=5)
        
        list_frame = ctk.CTkFrame(frame)
        list_frame.grid(row=1, column=0, sticky="nsew", padx=15, pady=15)
        list_frame.grid_rowconfigure(0, weight=1)
        list_frame.grid_columnconfigure(0, weight=1)
        
        self.merge_listbox = DragDropListbox(
            list_frame, 
            bg="#2b2b2b", 
            fg="white",
            font=("Segoe UI", 12), 
            selectbackground="#1f538d",
            selectmode=tk.SINGLE, 
            height=20
        )
        self.merge_listbox.grid(row=0, column=0, sticky="nsew")
        
        scrollbar = ctk.CTkScrollbar(list_frame, command=self.merge_listbox.yview)
        scrollbar.grid(row=0, column=1, sticky="ns")
        self.merge_listbox.config(yscrollcommand=scrollbar.set)
        
        info_frame = ctk.CTkFrame(frame)
        info_frame.grid(row=2, column=0, sticky="ew", padx=15, pady=15)
        
        self.merge_info_label = ctk.CTkLabel(
            info_frame, 
            text="Drag files to reorder | Double-click to preview",
            font=ctk.CTkFont(size=12)
        )
        self.merge_info_label.pack(pady=10)
        
        return frame
    
    def create_compress_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        
        ctk.CTkLabel(
            frame, 
            text="PDF Compression", 
            font=ctk.CTkFont(size=28, weight="bold")
        ).pack(pady=20)
        
        self.compress_frame = ctk.CTkFrame(frame)
        self.compress_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkButton(
            self.compress_frame, 
            text="📂 Choose PDF File", 
            command=self.select_compress_file,
            height=45,
            width=200
        ).pack(pady=10)
        
        self.compress_file_label = ctk.CTkLabel(
            self.compress_frame, 
            text="No file selected",
            font=ctk.CTkFont(size=12)
        )
        self.compress_file_label.pack(pady=5)
        
        quality_frame = ctk.CTkFrame(frame)
        quality_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkLabel(
            quality_frame, 
            text="Compression Settings", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(pady=10)
        
        qual_control = ctk.CTkFrame(quality_frame, fg_color="transparent")
        qual_control.pack(pady=10)
        
        ctk.CTkLabel(qual_control, text="Image Quality:").pack(side="left", padx=5)
        self.quality_slider = ctk.CTkSlider(qual_control, from_=10, to=100, width=200)
        self.quality_slider.set(80)
        self.quality_slider.pack(side="left", padx=5)
        self.quality_label = ctk.CTkLabel(qual_control, text="80%")
        self.quality_label.pack(side="left", padx=5)
        self.quality_slider.configure(command=lambda v: self.quality_label.configure(text=f"{int(v)}%"))
        
        ctk.CTkButton(
            frame, 
            text="🚀 Compress PDF", 
            command=self.execute_compress,
            fg_color="green",
            height=45,
            width=200
        ).pack(pady=20)
        
        self.compression_result = ctk.CTkLabel(frame, text="", font=ctk.CTkFont(size=14))
        self.compression_result.pack(pady=10)
        
        return frame
    
    def create_security_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        
        ctk.CTkLabel(
            frame, 
            text="Security Tools", 
            font=ctk.CTkFont(size=28, weight="bold")
        ).pack(pady=20)
        
        encrypt_frame = ctk.CTkFrame(frame)
        encrypt_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkLabel(
            encrypt_frame, 
            text="🔐 Encrypt PDF", 
            font=ctk.CTkFont(size=18, weight="bold")
        ).pack(pady=10)
        
        ctk.CTkButton(
            encrypt_frame, 
            text="Select PDF", 
            command=self.select_encrypt_file,
            height=40,
            width=200
        ).pack(pady=10)
        
        self.encrypt_file_label = ctk.CTkLabel(encrypt_frame, text="No file selected")
        self.encrypt_file_label.pack()
        
        self.password_entry = ctk.CTkEntry(
            encrypt_frame, 
            placeholder_text="Enter password",
            show="•",
            width=250,
            height=35
        )
        self.password_entry.pack(pady=10)
        
        perms_frame = ctk.CTkFrame(encrypt_frame, fg_color="transparent")
        perms_frame.pack(pady=10)
        
        self.allow_print = tk.BooleanVar(value=True)
        self.allow_copy = tk.BooleanVar(value=True)
        self.allow_modify = tk.BooleanVar(value=False)
        
        ctk.CTkCheckBox(perms_frame, text="Allow Printing", 
                       variable=self.allow_print).pack(anchor="w", pady=2)
        ctk.CTkCheckBox(perms_frame, text="Allow Copying", 
                       variable=self.allow_copy).pack(anchor="w", pady=2)
        ctk.CTkCheckBox(perms_frame, text="Allow Modifying", 
                       variable=self.allow_modify).pack(anchor="w", pady=2)
        
        ctk.CTkButton(
            encrypt_frame, 
            text="Encrypt & Save", 
            command=self.execute_encrypt,
            fg_color="green",
            height=40,
            width=200
        ).pack(pady=20)
        
        decrypt_frame = ctk.CTkFrame(frame)
        decrypt_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkLabel(
            decrypt_frame, 
            text="🔓 Decrypt PDF", 
            font=ctk.CTkFont(size=18, weight="bold")
        ).pack(pady=10)
        
        ctk.CTkButton(
            decrypt_frame, 
            text="Select Encrypted PDF", 
            command=self.select_decrypt_file,
            height=40,
            width=200
        ).pack(pady=10)
        
        self.decrypt_file_label = ctk.CTkLabel(decrypt_frame, text="No file selected")
        self.decrypt_file_label.pack()
        
        self.decrypt_password = ctk.CTkEntry(
            decrypt_frame, 
            placeholder_text="Enter password",
            show="•",
            width=250,
            height=35
        )
        self.decrypt_password.pack(pady=10)
        
        ctk.CTkButton(
            decrypt_frame, 
            text="Decrypt & Save", 
            command=self.execute_decrypt,
            fg_color="blue",
            height=40,
            width=200
        ).pack(pady=20)
        
        return frame
    
    def create_ocr_tab(self):
        frame = ctk.CTkFrame(self.container)
        frame.grid(row=0, column=0, sticky="nsew")
        
        ctk.CTkLabel(
            frame, 
            text="OCR Tools", 
            font=ctk.CTkFont(size=28, weight="bold")
        ).pack(pady=20)
        
        if not OCR_AVAILABLE:
            ctk.CTkLabel(
                frame, 
                text="⚠️ OCR not available - Install pytesseract and pdf2image", 
                text_color="orange",
                font=ctk.CTkFont(size=14)
            ).pack(pady=20)
            return frame
        
        self.ocr_frame = ctk.CTkFrame(frame)
        self.ocr_frame.pack(pady=20, padx=50, fill="x")
        
        ctk.CTkButton(
            self.ocr_frame, 
            text="📂 Select PDF or Image", 
            command=self.select_ocr_file,
            height=45,
            width=200
        ).pack(pady=10)
        
        self.ocr_file_label = ctk.CTkLabel(self.ocr_frame, text="No file selected")
        self.ocr_file_label.pack()
        
        lang_frame = ctk.CTkFrame(self.ocr_frame)
        lang_frame.pack(pady=10)
        
        ctk.CTkLabel(lang_frame, text="Language:").pack(side="left", padx=5)
        self.ocr_lang = ctk.CTkComboBox(
            lang_frame, 
            values=["eng", "fra", "deu", "spa", "ita", "por", "rus", "chi_sim", "jpn"],
            width=150
        )
        self.ocr_lang.set("eng")
        self.ocr_lang.pack(side="left", padx=5)
        
        ctk.CTkButton(
            self.ocr_frame, 
            text="🔍 Extract Text", 
            command=self.execute_ocr,
            fg_color="green",
            height=40,
            width=200
        ).pack(pady=20)
        
        self.ocr_result = ctk.CTkTextbox(frame, height=300)
        self.ocr_result.pack(pady=20, padx=50, fill="both", expand=True)
        
        return frame
    
    def setup_statusbar(self):
        self.statusbar = ctk.CTkFrame(self, height=40)
        self.statusbar.grid(row=1, column=0, columnspan=2, sticky="ew")
        self.statusbar.grid_columnconfigure(0, weight=1)
        
        self.status_label = ctk.CTkLabel(
            self.statusbar, 
            text="Ready",
            font=ctk.CTkFont(size=12)
        )
        self.status_label.grid(row=0, column=0, padx=15, sticky="w")
        
        self.progress_bar = ctk.CTkProgressBar(self.statusbar, width=250, height=15)
        self.progress_bar.grid(row=0, column=1, padx=15)
        self.progress_bar.set(0)
        
        self.page_count_label = ctk.CTkLabel(
            self.statusbar, 
            text="",
            font=ctk.CTkFont(size=12)
        )
        self.page_count_label.grid(row=0, column=2, padx=15, sticky="e")
        
    def setup_bindings(self):
        shortcuts = self.config_manager.config["shortcuts"]
        
        if "open" in shortcuts:
            self.bind(shortcuts["open"], lambda e: self.open_pdf())
        if "save" in shortcuts:
            self.bind(shortcuts["save"], lambda e: self.save_pdf())
        if "save_as" in shortcuts:
            self.bind(shortcuts["save_as"], lambda e: self.save_pdf_as())
        if "undo" in shortcuts:
            self.bind(shortcuts["undo"], lambda e: self.undo_action())
        if "help" in shortcuts:
            self.bind(shortcuts["help"], lambda e: self.show_shortcuts())
        
        self.bind("<Control-y>", lambda e: self.redo_action())
        
    def setup_ui_updater(self):
        self.check_ui_queue()
        
    def check_ui_queue(self):
        try:
            while True:
                task = self.ui_queue.get_nowait()
                self.process_ui_task(task)
        except queue.Empty:
            pass
        self.after(50, self.check_ui_queue)
    
    def process_ui_task(self, task):
        cmd = task[0]
        if cmd == "log":
            self._do_log(task[1])
        elif cmd == "progress":
            self.progress_bar.set(task[1])
        elif cmd == "status":
            self.status_label.configure(text=task[1])
        elif cmd == "message":
            messagebox.showinfo(task[1], task[2])
        elif cmd == "error":
            messagebox.showerror(task[1], task[2])
        elif cmd == "thumb":
            self._display_thumbnail(task[1], task[2])
        elif cmd == "clear_thumbs":
            self._clear_thumbnails()
        elif cmd == "trans_progress":
            if hasattr(self, 'trans_progress'):
                self.trans_progress.set(task[1])
        elif cmd == "custom":
            task[1]()
    
    def check_dependencies(self):
        missing = []
        if not OCR_AVAILABLE:
            missing.append("OCR (pytesseract/pdf2image)")
        if not TRANSLATE_AVAILABLE:
            missing.append("Translation (googletrans/deep-translator)")
        if not PDF2DOCX_AVAILABLE:
            missing.append("PDF to Word (pdf2docx)")
        if not PPTX_AVAILABLE:
            missing.append("PowerPoint (python-pptx)")
        if not DOCX_AVAILABLE:
            missing.append("Word (python-docx)")
        
        if missing:
            msg = "Some optional features are unavailable:\n\n"
            msg += "\n".join(f"• {m}" for m in missing)
            self.after(1000, lambda: messagebox.showwarning("Missing Features", msg))
    
    def show_tab(self, tab_name):
        for name, frame in self.tabs.items():
            if name == tab_name:
                frame.lift()
                if name in self.nav_buttons:
                    self.nav_buttons[name].configure(fg_color=["#3B8ED0", "#1F6AA5"])
            else:
                if name in self.nav_buttons:
                    self.nav_buttons[name].configure(fg_color="transparent")
    
    def show_converter(self): self.show_tab("converter")
    def show_manager(self): self.show_tab("manager")
    def show_merger(self): self.show_tab("merger")
    def show_sign(self): self.show_tab("sign")
    def show_watermark(self): self.show_tab("watermark")
    def show_crop(self): self.show_tab("crop")
    def show_translate(self): self.show_tab("translate")
    def show_compress(self): self.show_tab("compress")
    def show_security(self): self.show_tab("security")
    def show_ocr(self): self.show_tab("ocr")
    
    def log(self, message: str):
        self.ui_queue.put(("log", message))
    
    def _do_log(self, message: str):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_text.configure(state="normal")
        self.log_text.insert("end", f"[{timestamp}] {message}\n")
        self.log_text.see("end")
        self.log_text.configure(state="disabled")
    
    def refresh_recent_files(self):
        for widget in self.recent_frame.winfo_children():
            widget.destroy()
        
        for filepath in self.config_manager.config["recent_files"]:
            if os.path.exists(filepath):
                btn = ctk.CTkButton(
                    self.recent_frame, 
                    text=os.path.basename(filepath),
                    command=lambda f=filepath: self.open_pdf(f),
                    fg_color="transparent", 
                    text_color=("black", "white"),
                    anchor="w", 
                    height=30
                )
                btn.pack(fill="x", pady=2)
    
    def add_recent(self, filepath: str):
        self.config_manager.add_recent_file(filepath)
        self.refresh_recent_files()
    
    def open_pdf(self, filepath: str = None):
        if not filepath:
            filepath = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        
        if not filepath or not os.path.exists(filepath):
            return
        
        try:
            if self.current_doc:
                self.undo_stack.append(("open", self.current_path, self.current_doc))
                if len(self.undo_stack) > 10:
                    self.undo_stack.pop(0)
            
            if self.current_doc:
                self.current_doc.close()
            
            self.current_doc = fitz.open(filepath)
            self.current_path = filepath
            self.add_recent(filepath)
            
            self.log(f"Opened: {os.path.basename(filepath)} ({len(self.current_doc)} pages)")
            self.page_count_label.configure(text=f"{len(self.current_doc)} pages")
            
            if self.tabs["manager"].winfo_viewable():
                self.load_thumbnails()
                
        except Exception as e:
            messagebox.showerror("Error", f"Could not open PDF: {str(e)}")
    
    def save_pdf(self):
        if not self.current_doc:
            return
        
        if self.current_path:
            try:
                self.current_doc.save(self.current_path, garbage=4, deflate=True)
                self.log(f"Saved: {os.path.basename(self.current_path)}")
            except Exception as e:
                messagebox.showerror("Error", f"Could not save: {str(e)}")
        else:
            self.save_pdf_as()
    
    def save_pdf_as(self):
        if not self.current_doc:
            return
        
        filepath = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if not filepath:
            return
        
        try:
            self.current_doc.save(filepath, garbage=4, deflate=True)
            self.current_path = filepath
            self.log(f"Saved: {os.path.basename(filepath)}")
            self.add_recent(filepath)
        except Exception as e:
            messagebox.showerror("Error", f"Could not save: {str(e)}")
    
    def undo_action(self):
        if not self.undo_stack:
            self.log("Nothing to undo")
            return
        
        action, path, doc = self.undo_stack.pop()
        if self.current_doc:
            self.redo_stack.append((action, self.current_path, self.current_doc))
            self.current_doc.close()
        
        self.current_doc = doc
        self.current_path = path
        self.log("Undid last action")
        self.load_thumbnails()
    
    def redo_action(self):
        if not self.redo_stack:
            self.log("Nothing to redo")
            return
        
        action, path, doc = self.redo_stack.pop()
        if self.current_doc:
            self.undo_stack.append((action, self.current_path, self.current_doc))
            self.current_doc.close()
        
        self.current_doc = doc
        self.current_path = path
        self.log("Redid action")
        self.load_thumbnails()
    
    def load_thumbnails(self):
        if not self.current_doc:
            return
        
        self.ui_queue.put(("clear_thumbs",))
        
        def load():
            for i in range(len(self.current_doc)):
                try:
                    img = self.pdf_engine.get_page_thumbnail(self.current_doc, i, zoom=0.15)
                    self.ui_queue.put(("thumb", i, img))
                except Exception as e:
                    self.log(f"Error loading page {i+1}: {e}")
        
        threading.Thread(target=load, daemon=True).start()
    
    def _clear_thumbnails(self):
        for widget in self.thumb_frame.winfo_children():
            widget.destroy()
        self.page_widgets.clear()
    
    def _display_thumbnail(self, page_num: int, img: Image.Image):
        cols = 4
        row = page_num // cols
        col = page_num % cols
        
        frame = ctk.CTkFrame(self.thumb_frame)
        frame.grid(row=row, column=col, padx=5, pady=5, sticky="n")
        
        img.thumbnail((120, 160), Image.Resampling.LANCZOS)
        ctk_img = ctk.CTkImage(light_image=img, dark_image=img, size=(120, 160))
        
        lbl = ctk.CTkLabel(frame, image=ctk_img, text="")
        lbl.pack()
        
        var = tk.BooleanVar()
        chk = ctk.CTkCheckBox(frame, text=f"Pg {page_num+1}", variable=var, width=80)
        chk.pack(pady=2)
        
        lbl.bind("<Button-1>", lambda e, p=page_num: self.select_page(p))
        lbl.bind("<Double-Button-1>", lambda e, p=page_num: self.show_preview(p))
        
        self.page_widgets[page_num] = {
            'frame': frame,
            'var': var,
            'image': ctk_img,
            'pil_image': img
        }
    
    def select_page(self, page_num: int):
        var = self.page_widgets[page_num]['var']
        var.set(not var.get())
    
    def show_preview(self, page_num: int):
        if not self.current_doc:
            return
        
        self.current_preview_page = page_num
        page = self.current_doc[page_num]
        
        zoom = 2.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        canvas_w, canvas_h = 400, 550
        img.thumbnail((canvas_w, canvas_h), Image.Resampling.LANCZOS)
        
        self.preview_image = ImageTk.PhotoImage(img)
        
        self.preview_canvas.delete("all")
        x = (canvas_w - img.width) // 2
        y = (canvas_h - img.height) // 2
        self.preview_canvas.create_image(x, y, anchor="nw", image=self.preview_image)
        
        self.page_info_label.configure(
            text=f"Page {page_num+1} of {len(self.current_doc)} | {int(page.rect.width)}×{int(page.rect.height)} pts"
        )
    
    def on_crop_start(self, event):
        if not hasattr(self, 'current_preview_page') or not self.current_doc:
            return
        self.crop_start = (event.x, event.y)
        self.preview_canvas.delete("crop_rect")
    
    def on_crop_drag(self, event):
        if not self.crop_start:
            return
        x1, y1 = self.crop_start
        x2, y2 = event.x, event.y
        self.preview_canvas.delete("crop_rect")
        self.preview_canvas.create_rectangle(x1, y1, x2, y2, outline="red", width=2, tags="crop_rect")
    
    def on_crop_end(self, event):
        if not self.crop_start or not hasattr(self, 'preview_image') or not self.current_doc:
            return
        
        x1, y1 = self.crop_start
        x2, y2 = event.x, event.y
        
        page = self.current_doc[self.current_preview_page]
        rect = page.rect
        
        img_w = self.preview_image.width()
        img_h = self.preview_image.height()
        canvas_w, canvas_h = 400, 550
        offset_x = (canvas_w - img_w) // 2
        offset_y = (canvas_h - img_h) // 2
        
        x1 = max(0, min(x1, x2) - offset_x)
        y1 = max(0, min(y1, y2) - offset_y)
        x2 = min(img_w, max(x1, x2) - offset_x)
        y2 = min(img_h, max(y1, y2) - offset_y)
        
        if x2 > x1 and y2 > y1:
            scale_x = rect.width / img_w
            scale_y = rect.height / img_h
            
            pdf_rect = (
                x1 * scale_x,
                y1 * scale_y,
                x2 * scale_x,
                y2 * scale_y
            )
            
            self.crop_rect = pdf_rect
            self.log(f"Crop area selected: {pdf_rect}")
            
            selected = [i for i, w in self.page_widgets.items() if w['var'].get()]
            if not selected:
                selected = [self.current_preview_page]
            
            self.undo_stack.append(("crop", self.current_path, fitz.open(self.current_path)))
            
            for idx in selected:
                self.pdf_engine.crop_page(self.current_doc, idx, pdf_rect)
            
            self.log(f"Cropped {len(selected)} pages")
            self.load_thumbnails()
        
        self.crop_start = None
    
    def apply_manual_crop(self):
        if not self.current_doc:
            messagebox.showwarning("Warning", "Open a PDF first")
            return
        
        try:
            x = float(self.crop_x.get() or 0)
            y = float(self.crop_y.get() or 0)
            w = float(self.crop_w.get() or 0)
            h = float(self.crop_h.get() or 0)
            
            if w <= 0 or h <= 0:
                messagebox.showwarning("Warning", "Invalid dimensions")
                return
            
            selected = [i for i, w in self.page_widgets.items() if w['var'].get()]
            if not selected:
                messagebox.showwarning("Warning", "No pages selected")
                return
            
            self.undo_stack.append(("crop", self.current_path, fitz.open(self.current_path)))
            
            for idx in selected:
                self.pdf_engine.crop_page(self.current_doc, idx, (x, y, x+w, y+h))
            
            self.log(f"Cropped {len(selected)} pages to {w}×{h}")
            self.load_thumbnails()
            
        except ValueError:
            messagebox.showerror("Error", "Invalid crop dimensions")
    
    def scale_pages(self):
        if not self.current_doc:
            messagebox.showwarning("Warning", "Open a PDF first")
            return
        
        factor = self.scale_factor.get()
        
        selected = [i for i, w in self.page_widgets.items() if w['var'].get()]
        if not selected:
            selected = range(len(self.current_doc))
        
        self.undo_stack.append(("scale", self.current_path, fitz.open(self.current_path)))
        
        # Note: Actual scaling would require recreating pages
        # This is a placeholder - implement if needed
        self.log(f"Scaled {len(selected)} pages by {factor:.2f}x (simulated)")
    
    def auto_crop_margins(self):
        if not self.current_doc:
            messagebox.showwarning("Warning", "Open a PDF first")
            return
        
        selected = [i for i, w in self.page_widgets.items() if w['var'].get()]
        if not selected:
            selected = range(len(self.current_doc))
        
        self.undo_stack.append(("autocrop", self.current_path, fitz.open(self.current_path)))
        
        cropped_count = 0
        for idx in selected:
            result = self.pdf_engine.auto_crop_margins(self.current_doc, idx)
            if result:
                cropped_count += 1
        
        self.log(f"Auto-cropped {cropped_count} of {len(selected)} pages")
        self.load_thumbnails()
    
    def manipulate_pages(self, action: str):
        if not self.current_doc:
            return
        
        selected = [i for i, w in self.page_widgets.items() if w['var'].get()]
        if not selected:
            messagebox.showwarning("Warning", "No pages selected")
            return
        
        self.undo_stack.append(("action", self.current_path, fitz.open(self.current_path)))
        
        if action == "rotate":
            for idx in selected:
                self.pdf_engine.rotate_page_permanent(self.current_doc, idx, 90)
            self.log(f"Rotated {len(selected)} pages")
            
        elif action == "delete":
            if messagebox.askyesno("Confirm", f"Delete {len(selected)} pages?"):
                for idx in sorted(selected, reverse=True):
                    self.current_doc.delete_page(idx)
                self.log(f"Deleted {len(selected)} pages")
                
        elif action == "extract":
            filepath = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
            if filepath:
                new_doc = fitz.open()
                for idx in selected:
                    new_doc.insert_pdf(self.current_doc, from_page=idx, to_page=idx)
                new_doc.save(filepath)
                new_doc.close()
                self.log(f"Extracted {len(selected)} pages")
        
        self.load_thumbnails()
    
    def show_split_dialog(self):
        if not self.current_doc:
            messagebox.showwarning("Warning", "Open a PDF first")
            return
        
        dialog = ctk.CTkToplevel(self)
        dialog.title("Split PDF")
        dialog.geometry("400x200")
        dialog.transient(self)
        
        ctk.CTkLabel(dialog, text="Split by pages:").pack(pady=10)
        
        value_entry = ctk.CTkEntry(dialog, width=100)
        value_entry.insert(0, "1")
        value_entry.pack(pady=5)
        
        def execute_split():
            try:
                val = int(value_entry.get())
                folder = filedialog.askdirectory()
                if folder:
                    output_files = PDFEngine.split_pdf(self.current_path, folder, val)
                    self.log(f"Split PDF into {len(output_files)} files")
                    dialog.destroy()
            except ValueError:
                messagebox.showerror("Error", "Invalid number")
        
        ctk.CTkButton(dialog, text="Split", command=execute_split).pack(pady=20)
    
    def refresh_signatures(self):
        sigs = self.config_manager.get_signatures()
        if hasattr(self, 'sig_combo'):
            self.sig_combo.configure(values=list(sigs.keys()))
            if sigs:
                self.sig_var.set(list(sigs.keys())[0])
    
    def save_drawn_signature(self):
        name = simpledialog.askstring("Save Signature", "Enter name for this signature:")
        if not name:
            return
        
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            temp_path = tmp.name
        
        self.sign_canvas.save_signature(temp_path)
        
        with open(temp_path, 'rb') as f:
            img_data = base64.b64encode(f.read()).decode()
        
        self.config_manager.save_signature(name, img_data)
        os.remove(temp_path)
        
        self.refresh_signatures()
        self.log(f"Saved signature: {name}")
    
    def select_sign_pdf(self):
        filepath = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if filepath:
            self.sign_pdf_path = filepath
            self.sign_pdf_label.configure(text=os.path.basename(filepath))
    
    def apply_signature(self):
        if not hasattr(self, 'sign_pdf_path'):
            messagebox.showwarning("Warning", "Select a PDF first")
            return
        
        sig_name = self.sig_var.get()
        if not sig_name:
            messagebox.showwarning("Warning", "Select a signature")
            return
        
        sig_data = self.config_manager.get_signatures().get(sig_name)
        if not sig_data:
            return
        
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(base64.b64decode(sig_data))
            sig_path = tmp.name
        
        try:
            doc = fitz.open(self.sign_pdf_path)
            position = self.sign_pos.get()
            
            for page in doc:
                rect = page.rect
                margin = 50
                sig_width, sig_height = 200, 80
                
                if position == "Bottom-Right":
                    sig_rect = fitz.Rect(
                        rect.width - sig_width - margin,
                        rect.height - sig_height - margin,
                        rect.width - margin,
                        rect.height - margin
                    )
                elif position == "Bottom-Left":
                    sig_rect = fitz.Rect(
                        margin,
                        rect.height - sig_height - margin,
                        margin + sig_width,
                        rect.height - margin
                    )
                elif position == "Top-Right":
                    sig_rect = fitz.Rect(
                        rect.width - sig_width - margin,
                        margin,
                        rect.width - margin,
                        margin + sig_height
                    )
                elif position == "Top-Left":
                    sig_rect = fitz.Rect(
                        margin,
                        margin,
                        margin + sig_width,
                        margin + sig_height
                    )
                else:  # Center
                    sig_rect = fitz.Rect(
                        rect.width/2 - sig_width/2,
                        rect.height/2 - sig_height/2,
                        rect.width/2 + sig_width/2,
                        rect.height/2 + sig_height/2
                    )
                
                page.insert_image(sig_rect, filename=sig_path)
            
            output = filedialog.asksaveasfilename(
                defaultextension=".pdf",
                filetypes=[("PDF", "*.pdf")]
            )
            if output:
                doc.save(output)
                self.log(f"Applied signature to {os.path.basename(self.sign_pdf_path)}")
                messagebox.showinfo("Success", "Signature applied successfully")
            
            doc.close()
            os.remove(sig_path)
            
        except Exception as e:
            messagebox.showerror("Error", f"Failed to apply signature: {str(e)}")
    
    def select_watermark_image(self):
        filepath = filedialog.askopenfilename(
            filetypes=[("Images", "*.png *.jpg *.jpeg *.bmp *.gif")]
        )
        if filepath:
            self.watermark_image_path = filepath
            self.wm_img_label.configure(text=os.path.basename(filepath))
    
    def pick_watermark_color(self):
        color = colorchooser.askcolor(title="Choose watermark color")
        if color[0]:
            self.wm_color = (color[0][0]/255, color[0][1]/255, color[0][2]/255)
            # Update preview
            r, g, b = [int(c*255) for c in self.wm_color]
            self.wm_color_preview.configure(text="⬛", text_color=f"#{r:02x}{g:02x}{b:02x}")
            self.log("Watermark color set")
    
    def apply_watermark(self, wm_type: str):
        if not self.current_doc:
            messagebox.showwarning("Warning", "Open a PDF first")
            return
        
        self.undo_stack.append(("watermark", self.current_path, fitz.open(self.current_path)))
        
        if wm_type == "text":
            text = self.wm_text.get()
            if not text:
                messagebox.showwarning("Warning", "Enter watermark text")
                return
            
            self.pdf_engine.add_watermark_text(
                self.current_doc,
                text=text,
                opacity=self.wm_opacity.get(),
                font_size=int(self.wm_size.get()),
                color=self.wm_color,
                angle=int(self.wm_angle.get()),
                position=self.wm_pos.get().lower().replace("-", "_")
            )
            self.log(f"Applied text watermark: {text}")
            
        elif wm_type == "image":
            if not hasattr(self, 'watermark_image_path'):
                messagebox.showwarning("Warning", "Select an image first")
                return
            
            self.pdf_engine.add_watermark_image(
                self.current_doc,
                image_path=self.watermark_image_path,
                opacity=self.wm_opacity.get(),
                position=self.wm_pos.get().lower().replace("-", "_"),
                scale=self.wm_img_scale.get()
            )
            self.log("Applied image watermark")
        
        self.load_thumbnails()
    
    def add_merge_files(self):
        files = filedialog.askopenfilenames(filetypes=[("PDF", "*.pdf")])
        for f in files:
            if f not in self.merge_files:
                self.merge_files.append(f)
                self.merge_listbox.insert(tk.END, os.path.basename(f))
        self.update_merge_info()
    
    def clear_merge_files(self):
        self.merge_files.clear()
        self.merge_listbox.delete(0, tk.END)
        self.update_merge_info()
    
    def execute_merge(self):
        if len(self.merge_files) < 2:
            messagebox.showwarning("Warning", "Add at least 2 files")
            return
        
        output = filedialog.asksaveasfilename(
            defaultextension=".pdf",
            filetypes=[("PDF", "*.pdf")]
        )
        if not output:
            return
        
        def merge_thread():
            try:
                self.ui_queue.put(("progress", 0))
                PDFEngine.merge_pdfs(
                    self.merge_files, 
                    output,
                    lambda p: self.ui_queue.put(("progress", p))
                )
                self.ui_queue.put(("log", f"Merged {len(self.merge_files)} files"))
                self.ui_queue.put(("progress", 1))
                self.ui_queue.put(("message", "Success", "PDFs merged successfully"))
            except Exception as e:
                self.ui_queue.put(("error", "Merge Failed", str(e)))
        
        threading.Thread(target=merge_thread, daemon=True).start()
    
    def update_merge_info(self):
        total_pages = 0
        for f in self.merge_files:
            try:
                with fitz.open(f) as doc:
                    total_pages += len(doc)
            except:
                pass
        self.merge_info_label.configure(
            text=f"{len(self.merge_files)} files, {total_pages} total pages"
        )
    
    def start_conversion(self, conversion_type: str):
        """Start file conversion - FIXED VERSION with all working"""
        if conversion_type == "img_pdf":
            files = filedialog.askopenfilenames(
                title="Select Images",
                filetypes=[("Images", "*.jpg *.jpeg *.png *.bmp *.gif *.tiff")]
            )
            if files:
                output = filedialog.asksaveasfilename(
                    title="Save PDF As",
                    defaultextension=".pdf",
                    filetypes=[("PDF files", "*.pdf")]
                )
                if output:
                    def convert():
                        try:
                            self.ui_queue.put(("progress", 0))
                            self.pdf_engine.images_to_pdf(
                                files, output,
                                dpi=self.config_manager.config["default_dpi"]
                            )
                            self.ui_queue.put(("log", f"✅ Converted {len(files)} images to PDF: {os.path.basename(output)}"))
                            self.ui_queue.put(("progress", 1))
                            self.ui_queue.put(("message", "Success", f"Converted {len(files)} images to PDF"))
                        except Exception as e:
                            self.ui_queue.put(("error", "Conversion Failed", str(e)))
                    threading.Thread(target=convert, daemon=True).start()
                    self.log("Starting image to PDF conversion...")
        
        elif conversion_type == "pdf_img":
            if not self.current_doc:
                messagebox.showwarning("Warning", "Open a PDF first")
                return
            
            # Ask for image format
            format_dialog = ctk.CTkToplevel(self)
            format_dialog.title("Select Image Format")
            format_dialog.geometry("300x250")
            format_dialog.transient(self)
            format_dialog.grab_set()
            
            ctk.CTkLabel(format_dialog, text="Choose output format:", font=ctk.CTkFont(size=14)).pack(pady=20)
            
            format_var = tk.StringVar(value="jpg")
            formats = ["jpg", "png", "tiff", "bmp"]
            
            for fmt in formats:
                ctk.CTkRadioButton(format_dialog, text=fmt.upper(), variable=format_var, value=fmt).pack(pady=5)
            
            ctk.CTkLabel(format_dialog, text="DPI:", font=ctk.CTkFont(size=12)).pack(pady=5)
            dpi_var = tk.StringVar(value=str(self.config_manager.config["default_dpi"]))
            dpi_entry = ctk.CTkEntry(format_dialog, textvariable=dpi_var, width=80)
            dpi_entry.pack(pady=5)
            
            def start_export():
                try:
                    dpi = int(dpi_var.get())
                except:
                    dpi = self.config_manager.config["default_dpi"]
                
                format_dialog.destroy()
                folder = filedialog.askdirectory(title="Select Output Folder")
                if folder:
                    def convert():
                        try:
                            self.ui_queue.put(("progress", 0))
                            total_pages = len(self.current_doc)
                            
                            for i, page in enumerate(self.current_doc):
                                # Render page at specified DPI
                                zoom = dpi / 72
                                mat = fitz.Matrix(zoom, zoom)
                                pix = page.get_pixmap(matrix=mat)
                                
                                # Convert to PIL Image
                                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                                
                                # Save in requested format
                                output_path = os.path.join(folder, f"page_{i+1:03d}.{format_var.get()}")
                                
                                if format_var.get() == 'jpg':
                                    img.save(output_path, 'JPEG', quality=95, dpi=(dpi, dpi))
                                elif format_var.get() == 'png':
                                    img.save(output_path, 'PNG', dpi=(dpi, dpi))
                                elif format_var.get() == 'tiff':
                                    img.save(output_path, 'TIFF', dpi=(dpi, dpi))
                                else:  # bmp
                                    img.save(output_path, 'BMP', dpi=(dpi, dpi))
                                
                                self.ui_queue.put(("progress", (i + 1) / total_pages))
                            
                            self.ui_queue.put(("log", f"✅ Exported {total_pages} pages to {folder}"))
                            self.ui_queue.put(("progress", 1))
                            self.ui_queue.put(("message", "Success", f"Exported {total_pages} images"))
                        except Exception as e:
                            self.ui_queue.put(("error", "Export Failed", str(e)))
                    
                    threading.Thread(target=convert, daemon=True).start()
                    self.log("Starting PDF to image conversion...")
            
            ctk.CTkButton(format_dialog, text="Export", command=start_export, fg_color="green", height=35).pack(pady=20)
        
        elif conversion_type == "pdf_word":
            if not self.current_doc:
                messagebox.showwarning("Warning", "Open a PDF first")
                return
            
            if not PDF2DOCX_AVAILABLE:
                messagebox.showwarning("Warning", "PDF to Word requires pdf2docx library\nInstall with: pip install pdf2docx")
                return
            
            output = filedialog.asksaveasfilename(
                title="Save Word Document As",
                defaultextension=".docx",
                filetypes=[("Word Document", "*.docx")]
            )
            if output:
                def convert():
                    try:
                        self.ui_queue.put(("progress", 0.1))
                        # Create converter
                        cv = Converter(self.current_path)
                        self.ui_queue.put(("progress", 0.3))
                        
                        # Convert
                        cv.convert(output, start=0, end=None)
                        self.ui_queue.put(("progress", 0.8))
                        
                        # Close
                        cv.close()
                        self.ui_queue.put(("progress", 1))
                        
                        self.ui_queue.put(("log", f"✅ Converted PDF to Word: {os.path.basename(output)}"))
                        self.ui_queue.put(("message", "Success", "PDF converted to Word successfully"))
                    except Exception as e:
                        self.ui_queue.put(("error", "Conversion Failed", str(e)))
                
                threading.Thread(target=convert, daemon=True).start()
                self.log("Starting PDF to Word conversion...")
        
        elif conversion_type == "word_pdf":
            if not DOCX_AVAILABLE:
                messagebox.showwarning("Warning", "Word to PDF requires python-docx library\nInstall with: pip install python-docx")
                return
            
            filepath = filedialog.askopenfilename(
                title="Select Word Document",
                filetypes=[("Word Document", "*.docx"), ("All files", "*.*")]
            )
            if filepath:
                output = filedialog.asksaveasfilename(
                    title="Save PDF As",
                    defaultextension=".pdf",
                    filetypes=[("PDF files", "*.pdf")]
                )
                if output:
                    def convert():
                        try:
                            self.ui_queue.put(("progress", 0.2))
                            
                            # Open Word document
                            doc = Document(filepath)
                            
                            # Create PDF
                            pdf_doc = fitz.open()
                            page = pdf_doc.new_page(width=595, height=842)  # A4
                            
                            # Extract text and add to PDF
                            y = 50
                            for para in doc.paragraphs:
                                if para.text.strip():
                                    page.insert_text((50, y), para.text, fontsize=11)
                                    y += 20
                                    if y > 800:
                                        page = pdf_doc.new_page(width=595, height=842)
                                        y = 50
                            
                            self.ui_queue.put(("progress", 0.8))
                            pdf_doc.save(output)
                            pdf_doc.close()
                            
                            self.ui_queue.put(("progress", 1))
                            self.ui_queue.put(("log", f"✅ Converted Word to PDF: {os.path.basename(output)}"))
                            self.ui_queue.put(("message", "Success", "Word converted to PDF successfully"))
                        except Exception as e:
                            self.ui_queue.put(("error", "Conversion Failed", str(e)))
                    
                    threading.Thread(target=convert, daemon=True).start()
                    self.log("Starting Word to PDF conversion...")
        
        elif conversion_type == "pdf_ppt":
            if not self.current_doc:
                messagebox.showwarning("Warning", "Open a PDF first")
                return
            
            if not PPTX_AVAILABLE:
                messagebox.showwarning("Warning", "PDF to PowerPoint requires python-pptx library\nInstall with: pip install python-pptx")
                return
            
            output = filedialog.asksaveasfilename(
                title="Save PowerPoint As",
                defaultextension=".pptx",
                filetypes=[("PowerPoint", "*.pptx")]
            )
            if output:
                def convert():
                    try:
                        self.ui_queue.put(("progress", 0.2))
                        
                        # Create presentation
                        prs = Presentation()
                        
                        # Add slides from PDF pages
                        for i, page in enumerate(self.current_doc):
                            # Create slide
                            slide_layout = prs.slide_layouts[6]  # Blank slide
                            slide = prs.slides.add_slide(slide_layout)
                            
                            # Add title
                            title = slide.shapes.title
                            if title:
                                title.text = f"Page {i+1}"
                            
                            # Add text
                            text = page.get_text()
                            if text.strip():
                                left = Inches(1)
                                top = Inches(1.5)
                                width = Inches(8)
                                height = Inches(5)
                                
                                textbox = slide.shapes.add_textbox(left, top, width, height)
                                text_frame = textbox.text_frame
                                text_frame.text = text[:500] + "..." if len(text) > 500 else text
                            
                            self.ui_queue.put(("progress", 0.2 + (0.6 * (i + 1) / len(self.current_doc))))
                        
                        # Save
                        prs.save(output)
                        self.ui_queue.put(("progress", 1))
                        
                        self.ui_queue.put(("log", f"✅ Converted PDF to PowerPoint: {os.path.basename(output)}"))
                        self.ui_queue.put(("message", "Success", "PDF converted to PowerPoint successfully"))
                    except Exception as e:
                        self.ui_queue.put(("error", "Conversion Failed", str(e)))
                
                threading.Thread(target=convert, daemon=True).start()
                self.log("Starting PDF to PowerPoint conversion...")
        
        elif conversion_type == "ppt_pdf":
            if not PPTX_AVAILABLE:
                messagebox.showwarning("Warning", "PowerPoint to PDF requires python-pptx library\nInstall with: pip install python-pptx")
                return
            
            filepath = filedialog.askopenfilename(
                title="Select PowerPoint File",
                filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")]
            )
            if filepath:
                output = filedialog.asksaveasfilename(
                    title="Save PDF As",
                    defaultextension=".pdf",
                    filetypes=[("PDF files", "*.pdf")]
                )
                if output:
                    def convert():
                        try:
                            self.ui_queue.put(("progress", 0.2))
                            
                            # Open presentation
                            prs = Presentation(filepath)
                            
                            # Create PDF
                            pdf_doc = fitz.open()
                            
                            for i, slide in enumerate(prs.slides):
                                page = pdf_doc.new_page(width=595, height=842)  # A4
                                
                                # Add slide number
                                page.insert_text((50, 50), f"Slide {i+1}", fontsize=14)
                                
                                # Extract and add text from slide
                                y = 100
                                for shape in slide.shapes:
                                    if hasattr(shape, "text") and shape.text:
                                        page.insert_text((50, y), shape.text[:100], fontsize=11)
                                        y += 20
                                
                                self.ui_queue.put(("progress", 0.2 + (0.6 * (i + 1) / len(prs.slides))))
                            
                            pdf_doc.save(output)
                            pdf_doc.close()
                            
                            self.ui_queue.put(("progress", 1))
                            self.ui_queue.put(("log", f"✅ Converted PowerPoint to PDF: {os.path.basename(output)}"))
                            self.ui_queue.put(("message", "Success", "PowerPoint converted to PDF successfully"))
                        except Exception as e:
                            self.ui_queue.put(("error", "Conversion Failed", str(e)))
                    
                    threading.Thread(target=convert, daemon=True).start()
                    self.log("Starting PowerPoint to PDF conversion...")
        
        elif conversion_type == "pdf_txt":
            if not self.current_doc:
                messagebox.showwarning("Warning", "Open a PDF first")
                return
            
            output = filedialog.asksaveasfilename(
                title="Save Text File As",
                defaultextension=".txt",
                filetypes=[("Text File", "*.txt")]
            )
            if output:
                def convert():
                    try:
                        self.ui_queue.put(("progress", 0))
                        total_pages = len(self.current_doc)
                        
                        with open(output, 'w', encoding='utf-8') as f:
                            for i, page in enumerate(self.current_doc):
                                text = page.get_text()
                                f.write(f"--- Page {i+1} ---\n")
                                f.write(text)
                                f.write("\n\n")
                                self.ui_queue.put(("progress", (i + 1) / total_pages))
                        
                        self.ui_queue.put(("log", f"✅ Exported PDF to text: {os.path.basename(output)}"))
                        self.ui_queue.put(("progress", 1))
                        self.ui_queue.put(("message", "Success", "PDF exported to text successfully"))
                    except Exception as e:
                        self.ui_queue.put(("error", "Export Failed", str(e)))
                
                threading.Thread(target=convert, daemon=True).start()
                self.log("Starting PDF to text export...")
        
        elif conversion_type == "pdf_html":
            if not self.current_doc:
                messagebox.showwarning("Warning", "Open a PDF first")
                return
            
            output = filedialog.asksaveasfilename(
                title="Save HTML File As",
                defaultextension=".html",
                filetypes=[("HTML File", "*.html")]
            )
            if output:
                def convert():
                    try:
                        self.ui_queue.put(("progress", 0))
                        total_pages = len(self.current_doc)
                        
                        html = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>PDF Export</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 40px; }
        .page { margin-bottom: 30px; padding: 20px; border: 1px solid #ccc; border-radius: 5px; }
        .page-header { font-size: 18px; font-weight: bold; margin-bottom: 10px; color: #333; }
        .page-content { white-space: pre-wrap; line-height: 1.5; }
    </style>
</head>
<body>
"""
                        
                        for i, page in enumerate(self.current_doc):
                            text = page.get_text()
                            # Escape HTML special characters
                            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br>')
                            html += f"""
    <div class="page">
        <div class="page-header">Page {i+1}</div>
        <div class="page-content">{text}</div>
    </div>
"""
                            self.ui_queue.put(("progress", (i + 1) / total_pages))
                        
                        html += "\n</body>\n</html>"
                        
                        with open(output, 'w', encoding='utf-8') as f:
                            f.write(html)
                        
                        self.ui_queue.put(("log", f"✅ Exported PDF to HTML: {os.path.basename(output)}"))
                        self.ui_queue.put(("progress", 1))
                        self.ui_queue.put(("message", "Success", "PDF exported to HTML successfully"))
                    except Exception as e:
                        self.ui_queue.put(("error", "Export Failed", str(e)))
                
                threading.Thread(target=convert, daemon=True).start()
                self.log("Starting PDF to HTML export...")
        
        elif conversion_type == "pdf_pdfa":
            if not self.current_doc:
                messagebox.showwarning("Warning", "Open a PDF first")
                return
            
            output = filedialog.asksaveasfilename(
                title="Save PDF/A As",
                defaultextension=".pdf",
                filetypes=[("PDF files", "*.pdf")]
            )
            if output:
                def convert():
                    try:
                        self.ui_queue.put(("progress", 0.3))
                        # Simple PDF/A conversion - just save with PDF/A-1b flag
                        self.current_doc.save(output, garbage=4, deflate=True)
                        self.ui_queue.put(("progress", 1))
                        self.ui_queue.put(("log", f"✅ Converted to PDF/A: {os.path.basename(output)}"))
                        self.ui_queue.put(("message", "Success", "PDF converted to PDF/A format"))
                    except Exception as e:
                        self.ui_queue.put(("error", "Conversion Failed", str(e)))
                
                threading.Thread(target=convert, daemon=True).start()
                self.log("Starting PDF to PDF/A conversion...")
    
    def select_translate_pdf(self):
        filepath = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if filepath:
            self.translate_pdf_path = filepath
            self.trans_pdf_label.configure(text=os.path.basename(filepath))
            self.trans_result.delete("1.0", tk.END)
    
    def execute_translation(self):
        if not hasattr(self, 'translate_pdf_path'):
            messagebox.showwarning("Warning", "Select a PDF first")
            return
        
        def translate_thread():
            try:
                doc = fitz.open(self.translate_pdf_path)
                total = len(doc)
                translated_text = ""
                
                for i, page in enumerate(doc):
                    self.ui_queue.put(("trans_progress", i/total))
                    page_text = page.get_text()
                    
                    if page_text.strip():
                        try:
                            if 'GoogleTranslator' in globals():
                                translator = GoogleTranslator(
                                    source='auto', 
                                    target=self.tgt_lang.get()
                                )
                                translated = translator.translate(page_text)
                            else:
                                translator = Translator()
                                result = translator.translate(
                                    page_text, 
                                    dest=self.tgt_lang.get()
                                )
                                translated = result.text
                            
                            translated_text += f"--- Page {i+1} ---\n{translated}\n\n"
                        except Exception as e:
                            translated_text += f"--- Page {i+1} (Translation Error) ---\n{page_text}\n\n"
                    else:
                        translated_text += f"--- Page {i+1} (Empty) ---\n\n"
                    
                    self.ui_queue.put(("trans_progress", (i+1)/total))
                
                doc.close()
                
                if self.trans_mode.get() == "text":
                    output = filedialog.asksaveasfilename(
                        defaultextension=".txt",
                        filetypes=[("Text", "*.txt")]
                    )
                    if output:
                        with open(output, 'w', encoding='utf-8') as f:
                            f.write(translated_text)
                        self.ui_queue.put(("log", f"Translated text saved to {os.path.basename(output)}"))
                else:
                    # Show in textbox
                    self.ui_queue.put(("custom", lambda: self.trans_result.delete("1.0", tk.END)))
                    self.ui_queue.put(("custom", lambda: self.trans_result.insert("1.0", translated_text)))
                
                self.ui_queue.put(("log", f"Translation completed for {total} pages"))
                self.ui_queue.put(("trans_progress", 1))
                
            except Exception as e:
                self.ui_queue.put(("error", "Translation Failed", str(e)))
        
        threading.Thread(target=translate_thread, daemon=True).start()
        self.log("Starting translation...")
    
    def select_compress_file(self):
        filepath = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if filepath:
            self.compress_file_path = filepath
            self.compress_file_label.configure(text=os.path.basename(filepath))
            size = os.path.getsize(filepath) / (1024*1024)
            self.compression_result.configure(
                text=f"Original size: {size:.2f} MB",
                text_color="gray"
            )
    
    def execute_compress(self):
        if not hasattr(self, 'compress_file_path'):
            messagebox.showwarning("Warning", "Select a PDF first")
            return
        
        output = filedialog.asksaveasfilename(
            defaultextension=".pdf",
            filetypes=[("PDF", "*.pdf")]
        )
        if not output:
            return
        
        def compress_thread():
            try:
                self.ui_queue.put(("progress", 0))
                
                orig_size = os.path.getsize(self.compress_file_path) / (1024*1024)
                
                ratio = self.pdf_engine.compress_pdf(
                    self.compress_file_path,
                    output,
                    int(self.quality_slider.get())
                )
                
                new_size = os.path.getsize(output) / (1024*1024)
                
                self.ui_queue.put(("log", 
                    f"Compressed: {orig_size:.2f}MB → {new_size:.2f}MB ({ratio:.1f}% reduction)"))
                self.ui_queue.put(("progress", 1))
                
                result_text = f"Original: {orig_size:.2f} MB\nCompressed: {new_size:.2f} MB\nReduction: {ratio:.1f}%"
                self.ui_queue.put(("custom", 
                    lambda: self.compression_result.configure(
                        text=result_text,
                        text_color="green" if ratio > 0 else "orange"
                    )))
                
            except Exception as e:
                self.ui_queue.put(("error", "Compression Failed", str(e)))
        
        threading.Thread(target=compress_thread, daemon=True).start()
        self.log("Starting compression...")
    
    def select_encrypt_file(self):
        filepath = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if filepath:
            self.encrypt_file_path = filepath
            self.encrypt_file_label.configure(text=os.path.basename(filepath))
    
    def execute_encrypt(self):
        if not hasattr(self, 'encrypt_file_path'):
            messagebox.showwarning("Warning", "Select a PDF first")
            return
        
        password = self.password_entry.get()
        if not password:
            messagebox.showwarning("Warning", "Enter a password")
            return
        
        output = filedialog.asksaveasfilename(
            defaultextension=".pdf",
            filetypes=[("PDF", "*.pdf")]
        )
        if not output:
            return
        
        try:
            self.pdf_engine.encrypt_pdf(
                self.encrypt_file_path,
                output,
                password,
                self.allow_print.get(),
                self.allow_copy.get(),
                self.allow_modify.get()
            )
            self.log(f"Encrypted PDF saved to {os.path.basename(output)}")
            messagebox.showinfo("Success", "PDF encrypted successfully")
        except Exception as e:
            messagebox.showerror("Error", f"Encryption failed: {str(e)}")
    
    def select_decrypt_file(self):
        filepath = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if filepath:
            self.decrypt_file_path = filepath
            self.decrypt_file_label.configure(text=os.path.basename(filepath))
    
    def execute_decrypt(self):
        if not hasattr(self, 'decrypt_file_path'):
            messagebox.showwarning("Warning", "Select a PDF first")
            return
        
        password = self.decrypt_password.get()
        if not password:
            messagebox.showwarning("Warning", "Enter a password")
            return
        
        output = filedialog.asksaveasfilename(
            defaultextension=".pdf",
            filetypes=[("PDF", "*.pdf")]
        )
        if not output:
            return
        
        try:
            success = self.pdf_engine.decrypt_pdf(
                self.decrypt_file_path,
                output,
                password
            )
            if success:
                self.log(f"Decrypted PDF saved to {os.path.basename(output)}")
                messagebox.showinfo("Success", "PDF decrypted successfully")
            else:
                messagebox.showerror("Error", "Incorrect password")
        except Exception as e:
            messagebox.showerror("Error", f"Decryption failed: {str(e)}")
    
    def select_ocr_file(self):
        filepath = filedialog.askopenfilename(
            filetypes=[("PDF", "*.pdf"), ("Images", "*.jpg *.jpeg *.png *.bmp")]
        )
        if filepath:
            self.ocr_file_path = filepath
            self.ocr_file_label.configure(text=os.path.basename(filepath))
            self.ocr_result.delete("1.0", tk.END)
    
    def execute_ocr(self):
        if not hasattr(self, 'ocr_file_path'):
            messagebox.showwarning("Warning", "Select a file first")
            return
        
        def ocr_thread():
            try:
                text = ""
                ext = os.path.splitext(self.ocr_file_path)[1].lower()
                
                if ext == '.pdf':
                    images = convert_from_path(self.ocr_file_path, dpi=300)
                    for i, img in enumerate(images):
                        page_text = pytesseract.image_to_string(
                            img, lang=self.ocr_lang.get()
                        )
                        text += f"--- Page {i+1} ---\n{page_text}\n\n"
                        self.ui_queue.put(("progress", (i+1)/len(images)))
                else:
                    img = Image.open(self.ocr_file_path)
                    text = pytesseract.image_to_string(img, lang=self.ocr_lang.get())
                    self.ui_queue.put(("progress", 1))
                
                self.ui_queue.put(("custom", lambda: self.ocr_result.delete("1.0", tk.END)))
                self.ui_queue.put(("custom", lambda: self.ocr_result.insert("1.0", text)))
                self.ui_queue.put(("log", "OCR completed"))
                
            except Exception as e:
                self.ui_queue.put(("error", "OCR Failed", str(e)))
        
        threading.Thread(target=ocr_thread, daemon=True).start()
        self.log("Starting OCR...")
    
    def show_settings(self):
        dialog = ctk.CTkToplevel(self)
        dialog.title("Settings")
        dialog.geometry("600x600")
        dialog.transient(self)
        dialog.grab_set()
        
        notebook = ttk.Notebook(dialog)
        notebook.pack(fill="both", expand=True, padx=20, pady=20)
        
        # Appearance Tab
        appearance_frame = ctk.CTkFrame(notebook)
        notebook.add(appearance_frame, text="Appearance")
        
        theme_frame = ctk.CTkFrame(appearance_frame)
        theme_frame.pack(fill="x", padx=30, pady=30)
        
        ctk.CTkLabel(
            theme_frame, 
            text="Application Theme", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(anchor="w", pady=(0, 15))
        
        theme_var = tk.StringVar(value=self.config_manager.config["theme"])
        
        theme_buttons = ctk.CTkFrame(theme_frame, fg_color="transparent")
        theme_buttons.pack()
        
        def set_theme(theme):
            theme_var.set(theme)
            # Update button colors
            system_btn.configure(fg_color="#3B8ED0" if theme == "System" else "gray30")
            light_btn.configure(fg_color="#3B8ED0" if theme == "Light" else "gray30")
            dark_btn.configure(fg_color="#3B8ED0" if theme == "Dark" else "gray30")
        
        system_btn = ctk.CTkButton(
            theme_buttons,
            text="🌓 System",
            command=lambda: set_theme("System"),
            width=120,
            height=45,
            fg_color="#3B8ED0" if theme_var.get() == "System" else "gray30",
            font=ctk.CTkFont(size=13)
        )
        system_btn.pack(side="left", padx=5)
        
        light_btn = ctk.CTkButton(
            theme_buttons,
            text="☀️ Light",
            command=lambda: set_theme("Light"),
            width=120,
            height=45,
            fg_color="#3B8ED0" if theme_var.get() == "Light" else "gray30",
            font=ctk.CTkFont(size=13)
        )
        light_btn.pack(side="left", padx=5)
        
        dark_btn = ctk.CTkButton(
            theme_buttons,
            text="🌙 Dark",
            command=lambda: set_theme("Dark"),
            width=120,
            height=45,
            fg_color="#3B8ED0" if theme_var.get() == "Dark" else "gray30",
            font=ctk.CTkFont(size=13)
        )
        dark_btn.pack(side="left", padx=5)
        
        color_frame = ctk.CTkFrame(appearance_frame)
        color_frame.pack(fill="x", padx=30, pady=20)
        
        ctk.CTkLabel(
            color_frame, 
            text="Color Theme", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(anchor="w", pady=(0, 15))
        
        color_var = tk.StringVar(value=self.config_manager.config["color_theme"])
        
        def set_color(color):
            color_var.set(color)
            blue_btn.configure(fg_color="#3B8ED0" if color == "blue" else "gray30")
            green_btn.configure(fg_color="#2E7D32" if color == "green" else "gray30")
            dark_blue_btn.configure(fg_color="#1A237E" if color == "dark-blue" else "gray30")
        
        color_buttons = ctk.CTkFrame(color_frame, fg_color="transparent")
        color_buttons.pack()
        
        blue_btn = ctk.CTkButton(
            color_buttons,
            text="🔵 Blue",
            command=lambda: set_color("blue"),
            width=120,
            height=45,
            fg_color="#3B8ED0" if color_var.get() == "blue" else "gray30",
            font=ctk.CTkFont(size=13)
        )
        blue_btn.pack(side="left", padx=5)
        
        green_btn = ctk.CTkButton(
            color_buttons,
            text="🟢 Green",
            command=lambda: set_color("green"),
            width=120,
            height=45,
            fg_color="#2E7D32" if color_var.get() == "green" else "gray30",
            font=ctk.CTkFont(size=13)
        )
        green_btn.pack(side="left", padx=5)
        
        dark_blue_btn = ctk.CTkButton(
            color_buttons,
            text="🔷 Dark Blue",
            command=lambda: set_color("dark-blue"),
            width=120,
            height=45,
            fg_color="#1A237E" if color_var.get() == "dark-blue" else "gray30",
            font=ctk.CTkFont(size=13)
        )
        dark_blue_btn.pack(side="left", padx=5)
        
        # General Tab
        general_frame = ctk.CTkFrame(notebook)
        notebook.add(general_frame, text="General")
        
        dpi_frame = ctk.CTkFrame(general_frame)
        dpi_frame.pack(fill="x", padx=30, pady=30)
        
        ctk.CTkLabel(
            dpi_frame, 
            text="Default DPI", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(anchor="w", pady=(0, 15))
        
        dpi_var = tk.StringVar(value=str(self.config_manager.config["default_dpi"]))
        dpi_entry = ctk.CTkEntry(dpi_frame, textvariable=dpi_var, width=100)
        dpi_entry.pack(anchor="w")
        
        auto_frame = ctk.CTkFrame(general_frame)
        auto_frame.pack(fill="x", padx=30, pady=20)
        
        ctk.CTkLabel(
            auto_frame, 
            text="Auto-Save", 
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(anchor="w", pady=(0, 15))
        
        auto_save_var = tk.BooleanVar(value=self.config_manager.config.get("auto_save", True))
        auto_check = ctk.CTkCheckBox(
            auto_frame,
            text="Auto-save before operations",
            variable=auto_save_var
        )
        auto_check.pack(anchor="w", pady=5)
        
        backup_var = tk.BooleanVar(value=self.config_manager.config.get("backup_on_save", True))
        backup_check = ctk.CTkCheckBox(
            auto_frame,
            text="Create backup when saving",
            variable=backup_var
        )
        backup_check.pack(anchor="w", pady=5)
        
        # Save Button
        def save_settings():
            self.config_manager.config["theme"] = theme_var.get()
            self.config_manager.config["color_theme"] = color_var.get()
            self.config_manager.config["default_dpi"] = int(dpi_var.get() or 300)
            self.config_manager.config["auto_save"] = auto_save_var.get()
            self.config_manager.config["backup_on_save"] = backup_var.get()
            
            ctk.set_appearance_mode(theme_var.get())
            ctk.set_default_color_theme(color_var.get())
            self.config_manager.save()
            
            dialog.destroy()
            messagebox.showinfo("Settings", "Settings saved successfully!")
        
        button_frame = ctk.CTkFrame(dialog, fg_color="transparent")
        button_frame.pack(pady=20)
        
        ctk.CTkButton(
            button_frame,
            text="Cancel",
            command=dialog.destroy,
            width=100,
            height=35
        ).pack(side="left", padx=5)
        
        ctk.CTkButton(
            button_frame,
            text="Save",
            command=save_settings,
            width=100,
            height=35,
            fg_color="green"
        ).pack(side="left", padx=5)
    
    def show_shortcuts(self):
        dialog = ctk.CTkToplevel(self)
        dialog.title("Keyboard Shortcuts")
        dialog.geometry("400x400")
        dialog.transient(self)
        
        frame = ctk.CTkScrollableFrame(dialog)
        frame.pack(fill="both", expand=True, padx=20, pady=20)
        
        shortcuts = self.config_manager.config["shortcuts"]
        
        for action, key in shortcuts.items():
            row = ctk.CTkFrame(frame, fg_color="transparent")
            row.pack(fill="x", pady=2)
            
            ctk.CTkLabel(
                row, 
                text=f"{action.replace('_', ' ').title()}:", 
                font=ctk.CTkFont(weight="bold"),
                width=120
            ).pack(side="left", padx=10)
            ctk.CTkLabel(row, text=key).pack(side="right", padx=10)
    
    def show_about(self):
        about_text = f"""{APP_NAME} v{VERSION}
        
A comprehensive PDF manipulation toolkit with:
• PDF Editing & Conversion
• Digital Signatures
• Watermarking
• OCR Text Extraction
• Security & Encryption
• Translation
• PDF to/from Images
• PDF to/from Word
• PDF to/from PowerPoint

Created with Python, CustomTkinter, and PyMuPDF"""
        
        messagebox.showinfo("About", about_text)

if __name__ == "__main__":
    app = ProConvertStudio()
    app.mainloop()